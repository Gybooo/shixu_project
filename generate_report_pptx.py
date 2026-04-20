"""
生成项目报告 PPT (v2)
修改要点:
1. 明亮色风格 (浅蓝 + 白 + 橙)
2. 每页有明确「图 N-M / 图表名称」
3. 参数/术语说明放到演讲者备注里 (不出现在 slide 上)
4. slide 上只保留: 标题 + 图 + 简短一句话观点
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT_DIR = r"E:\timesfm_project\output"
PPT_PATH = os.path.join(OUTPUT_DIR, "振动预测项目报告_v2.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# 明亮色调
PRIMARY = RGBColor(0x4A, 0x90, 0xE2)    # 天蓝
ACCENT = RGBColor(0xF5, 0xA6, 0x23)     # 暖橙
SUCCESS = RGBColor(0x7E, 0xD3, 0x21)    # 草绿
BG = RGBColor(0xFC, 0xFD, 0xFE)          # 近白
LIGHT = RGBColor(0xE8, 0xF0, 0xFA)       # 淡蓝
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)   # 深灰
TEXT_MID = RGBColor(0x66, 0x6E, 0x7A)    # 中灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_font(run, name="微软雅黑", size=14, bold=False, color=TEXT_DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for e in rPr.findall(qn("a:ea")):
        rPr.remove(e)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)


def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.text = ""
        run = p.add_run()
        run.text = line
        set_font(run, name=font_name, size=size, bold=bold, color=color)
    return tb


def set_slide_bg(slide, color=BG):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, note_text):
    """添加演讲者备注"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = ""
    lines = note_text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line


def slide_header(slide, chart_no, chart_name, subtitle=None):
    """
    顶部: 浅蓝色条带, 显示:
      [图 N-M]  图表名称        副标题 (若有)
    """
    # 顶部细条
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=PRIMARY)

    # 图表编号徽章
    add_rect(slide, Inches(0.4), Inches(0.3), Inches(1.8), Inches(0.45),
             fill=ACCENT)
    add_text(slide, Inches(0.4), Inches(0.32), Inches(1.8), Inches(0.4),
             chart_no, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 图表名称
    add_text(slide, Inches(2.4), Inches(0.3), Inches(10.5), Inches(0.5),
             chart_name, size=20, bold=True, color=PRIMARY)

    # 副标题
    if subtitle:
        add_text(slide, Inches(2.4), Inches(0.78), Inches(10.5), Inches(0.3),
                 subtitle, size=11, color=TEXT_MID)


def add_figure_area(slide, img_path, top=Inches(1.35), max_height=Inches(5.2)):
    """居中放图, 自适应尺寸"""
    if not os.path.exists(img_path):
        return
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    max_w = SLIDE_W - Inches(0.8)
    ratio = iw / ih
    if max_height * ratio > max_w:
        w = max_w
        h = w / ratio
    else:
        h = max_height
        w = h * ratio
    left = (SLIDE_W - w) // 2
    slide.shapes.add_picture(img_path, left, top, width=w, height=h)
    return top + h


def add_key_point(slide, text, top=None):
    """底部单行重点 (明亮色背景)"""
    if top is None:
        top = Inches(6.7)
    add_rect(slide, Inches(0.4), top, SLIDE_W - Inches(0.8), Inches(0.55),
             fill=LIGHT)
    add_text(slide, Inches(0.6), top + Inches(0.05),
             SLIDE_W - Inches(1.2), Inches(0.45),
             "▍核心观点  " + text, size=13, bold=True, color=PRIMARY,
             anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_no, total):
    """页脚"""
    add_text(slide, Inches(11.8), SLIDE_H - Inches(0.3),
             Inches(1.5), Inches(0.3),
             f"{page_no} / {total}",
             size=10, color=TEXT_MID, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.3),
             Inches(11), Inches(0.3),
             "MPB_01 振动预测与预防性维护研究",
             size=10, color=TEXT_MID)


# ============================================================
# 每页内容: (chart_no, title, subtitle, image_file, key_point, notes)
# ============================================================

SLIDES_SPEC = [
    # (特殊: 封面 None)
    ("COVER", "MPB_01 设备振动预测与预防性维护研究",
     "LSTM vs TimesFM · 信号分解 · 多参数对比 · 泛化性验证",
     None, None,
     "本次汇报覆盖 4 大主题:\n"
     "  1. 11 属性原始数据的可视化与统计分析\n"
     "  2. LSTM 与 TimesFM 两种算法在不同参数下的对比\n"
     "  3. 为什么预测结果会退化为直线——诊断分析\n"
     "  4. 信号分解+预测窗口优化后的有效方案\n"
     "数据来源: MPB01_6.11 08_18.csv (InfluxDB 导出格式)\n"
     "数据规模: 30 万行, 11 属性, ~2.6 小时连续监测"),

    ("AGENDA", "报告目录",
     "本次汇报分 5 个章节",
     None, None,
     "章节 1: 数据概览与可视化 (slide 3-7)\n"
     "章节 2: 两种算法多参数对比 (slide 8-10)\n"
     "章节 3: 预测直线问题诊断 (slide 11)\n"
     "章节 4: 信号分解与窗口优化 (slide 12-16)\n"
     "章节 5: 多属性泛化测试与结论 (slide 17-20)"),

    ("图 1-1", "11 属性分量级时序对照",
     "按数值量级分 3 组, 橙色带标注 18 次疑似停机",
     "属性总览_1_分量级时序.png",
     "Magnitude 存在 18 次骤降事件, 对应设备停机或传感器离线",
     "图表内容:\n"
     "  - 【组A】Magnitude (合成幅值, 0-1 量级)\n"
     "  - 【组B】aRMS/vRMS 系列 (均方根, 0-5 量级)\n"
     "  - 【组C】ShockX/Y/Z (三轴冲击, ±1500 量级)\n"
     "  - 橙色半透明带 = 检测到的 18 次疑似异常事件\n"
     "\n"
     "参数/术语:\n"
     "  - RMS (Root Mean Square) 均方根: 衡量信号能量\n"
     "  - Magnitude: 三轴合成后的总幅值\n"
     "  - aRMS (acceleration RMS): 加速度均方根, 单位 g\n"
     "  - vRMS (velocity RMS): 速度均方根, 单位 mm/s\n"
     "  - Shock: 冲击峰值, 反映瞬时撞击强度\n"
     "  - 重采样频率 = 1 秒, 每个子图有 8960 个点"),

    ("图 1-2", "11 属性独立时序图",
     "每个属性一个面板, 便于观察各自量级和波动",
     "属性总览_2_独立子图.png",
     "Shock 类是尖峰脉冲信号, RMS 类是连续起伏, 物理本质不同",
     "图表内容:\n"
     "  - 11 个属性各自独立面板, 以便突出各自特点\n"
     "  - 每个子图标题含数据范围和标准差\n"
     "  - 红色虚线为均值线, 便于观察偏离\n"
     "\n"
     "参数/术语:\n"
     "  - 均值 (mean): 所有值的算术平均\n"
     "  - STD (Standard Deviation) 标准差: 衡量偏离均值的程度\n"
     "  - 范围 [min, max]: 数据上下限\n"
     "\n"
     "观察:\n"
     "  - ShockX/Y/Z: 密集尖峰(脉冲型)\n"
     "  - aRMS/vRMS: 随时间起伏(连续型)\n"
     "  - Magnitude: 平稳段 + 多次骤降(停机)"),

    ("图 1-3", "多属性统计概览与相关性",
     "log 轴量级对比 + 变异系数 + 11×11 相关矩阵",
     "属性总览_3_统计与相关性.png",
     "Magnitude/aRMS/vRMS 本质同一物理量 (相关 0.86~0.98), Shock 独立",
     "图表内容:\n"
     "  左上: 各属性绝对值均值 log 柱状图 (揭示跨量级差异)\n"
     "  右上: 变异系数 CV, ShockY=1.66 最难预测\n"
     "  下方: 11×11 皮尔逊相关系数热力图\n"
     "\n"
     "参数/术语:\n"
     "  - CV (变异系数) = STD / |均值|, 衡量相对波动\n"
     "    CV > 1   = 高难度 (红色线)\n"
     "    CV > 0.5 = 中难度 (橙色线)\n"
     "  - 皮尔逊相关系数 r ∈ [-1, 1]:\n"
     "    |r| > 0.7 强相关\n"
     "    |r| < 0.3 几乎无关\n"
     "\n"
     "关键发现:\n"
     "  - Magnitude 与 aRMS/vRMS 相关 0.86-0.98 (本质同类)\n"
     "  - Shock 与其他独立 (相关 ~0.05)\n"
     "  - ShockX/Z 和 ShockY/Z 负相关 (机械对称性)"),

    ("图 1-4", "异常事件放大 — Magnitude 骤降",
     "18 次事件的前 4 个详细波形",
     "属性总览_4_异常事件放大.png",
     "典型事件: 振动从 ~0.86 骤降到 ~0, 持续 3-14 秒后恢复",
     "图表内容:\n"
     "  展示 18 次疑似异常事件中的前 4 个\n"
     "  每个子图放大显示骤降前后 30 秒\n"
     "  橙色阴影 = 异常持续区间\n"
     "  灰色虚线 = 正常中位数\n"
     "\n"
     "参数/术语:\n"
     "  - z_threshold = 3.0 (异常判定阈值, 越低越敏感)\n"
     "  - min_duration = 3 秒 (最短持续, 避免瞬时误判)\n"
     "  - 判定规则: 值 < 中位数 - 3.0 * STD * 0.3\n"
     "\n"
     "业务含义:\n"
     "  - 骤降到 0 很可能是设备停机或传感器离线\n"
     "  - 对预测评估的启示: 测试集若含异常段会干扰评估"),

    ("图 2-1", "LSTM vs TimesFM 分维度参数对比",
     "在 Magnitude 上测试 3 个维度 × 7 组实验的 MAE",
     "对比_1_分维度柱状图.png",
     "TimesFM 在绝大多数配置下都优于 LSTM",
     "图表内容:\n"
     "  横向 3 个面板, 分别展示 3 个参数维度:\n"
     "  - 维度1 重采样频率: 1s / 5s / 30s\n"
     "  - 维度2 预测长度 horizon: 8 / 32 / 64\n"
     "  - 维度3 上下文长度 context: 128 / 256 / 512\n"
     "  红柱 = LSTM, 绿柱 = TimesFM, 越低越好\n"
     "\n"
     "参数/术语:\n"
     "  - MAE (Mean Absolute Error) 平均绝对误差\n"
     "    = mean(|预测值 - 真实值|), 单位与信号同\n"
     "  - horizon: 一次预测未来几步\n"
     "  - context: 模型看多长的历史数据\n"
     "\n"
     "结果:\n"
     "  - 5s 重采样比 1s 稍好 (降噪)\n"
     "  - horizon 越小 MAE 越小 (未来越近越好预测)\n"
     "  - TimesFM 对 context 不敏感, LSTM 对 context 敏感"),

    ("图 2-2", "6 组参数下的预测曲线样例",
     "发现问题: 预测普遍退化为近似直线",
     "对比_2_预测曲线样例.png",
     "两种模型的预测曲线都偏平, 触发后续诊断",
     "图表内容:\n"
     "  6 个子图对应 6 组不同参数实验\n"
     "  深蓝 = 真实值, 红虚线 = LSTM, 绿点划线 = TimesFM\n"
     "  可见: 即便是最好的 horizon=8 (实验 B1), 模型曲线\n"
     "  仍显著平于真实波动\n"
     "\n"
     "参数/术语:\n"
     "  - 真实值: 测试集的实际振动值\n"
     "  - LSTM: Encoder-Decoder + 直接多步输出, 在本地训练\n"
     "  - TimesFM: Google 预训练基础模型 2.5-200M\n"
     "\n"
     "启示:\n"
     "  - 表象 MAE 不差, 但视觉跟随效果不佳\n"
     "  - 需要后续诊断找出根本原因"),

    ("图 2-3", "算法胜率与 MAE 相对差距汇总",
     "TimesFM 在 5/6 组中胜出 (83%), 平均领先 41%",
     "对比_3_胜率汇总.png",
     "TimesFM 零样本即用, LSTM 需训练但效果不稳定",
     "图表内容:\n"
     "  左图: 胜率饼图, TimesFM 83% vs LSTM 17%\n"
     "  右图: 每组实验的 MAE 相对差距横向条形\n"
     "\n"
     "参数/术语:\n"
     "  - MAE 差距% = (差的/好的 - 1) × 100%\n"
     "    例如 LSTM=0.112, TFM=0.079, 差距 = 42.9%\n"
     "\n"
     "时间成本:\n"
     "  - TimesFM 总耗时 71 秒 (仅推理)\n"
     "  - LSTM 总耗时 137 秒 (训练+推理)\n"
     "  - 越大模型越大数据越能体现出 TimesFM 优势"),

    ("图 3-1", "诊断: 为什么预测结果是直线?",
     "自相关 + 朴素基线 + 功率谱三维度分析",
     "诊断_1_可预测性分析.png",
     "1 秒尺度接近高频噪声, 最优预测本就是均值线",
     "图表内容:\n"
     "  左上: 不同重采样尺度下的自相关函数 (ACF)\n"
     "  右上: 1s 信号功率谱 PSD\n"
     "  中间: 朴素基线 vs LSTM/TimesFM MAE 对比柱状\n"
     "  左下: 各尺度的相对误差\n"
     "  右下: 100 秒真实振动切片 (高频抖动)\n"
     "\n"
     "参数/术语:\n"
     "  - ACF (AutoCorrelation Function) 自相关函数\n"
     "    r(lag=k) = x[t] 与 x[t+k] 的相关性\n"
     "    r > 0.3 才有预测价值 (红色阈值线)\n"
     "  - 朴素基线: 不用模型的简单方法\n"
     "    例 朴素-最后值: 未来每步都等于当前最后值\n"
     "  - PSD (Power Spectral Density) 功率谱\n"
     "\n"
     "关键发现:\n"
     "  - 1s lag-1 自相关 0.87 (相邻秒高度相关)\n"
     "  - 但 lag-60 只有 0.09 (远期几乎独立)\n"
     "  - 朴素-最后值 MAE=0.069, 比 LSTM (0.112) 还好\n"
     "  - LSTM 训练结果不佳, 还不如照搬上一秒"),

    ("图 4-1", "方案 B - 信号分解对比",
     "把 Magnitude 拆成「趋势 + 噪声」, 测试 3 种方法",
     "分解_1_信号分解对比.png",
     "SavGol 保留趋势形状最好, Butter 最平滑但损失细节",
     "图表内容:\n"
     "  第一行: 原始信号(灰) + 3 种趋势(彩色)叠加\n"
     "  后 3 行: 每种分解的「趋势(主)+残差(副轴)」\n"
     "\n"
     "参数/术语:\n"
     "  - 移动均值 (Moving Average): 窗口内取平均\n"
     "    window = 30 秒\n"
     "  - Savitzky-Golay 滤波: 局部多项式拟合\n"
     "    window_length = 61, polyorder = 3\n"
     "  - Butterworth 低通滤波: 频域滤波\n"
     "    cutoff_hz = 0.05, order = 4\n"
     "  - 趋势占比 = Var(趋势) / Var(原信号)\n"
     "\n"
     "结果:\n"
     "  - 移动均值 70.5% / SavGol 71.9% / Butter 84.1%\n"
     "  - Butter 趋势占比最高, 但过于平滑"),

    ("图 4-2", "方案 B - 分解后可预测性跃升",
     "ACF 从 0.87 提升到 1.00, MAE 下降 47%",
     "分解_2_可预测性对比.png",
     "分解后的趋势信号真正可预测, 模型不再只能预测均值",
     "图表内容:\n"
     "  左图: ACF 对比 (lag-1 蓝, lag-30 紫)\n"
     "  右图: MAE 对比 (朴素灰, LSTM 红, TimesFM 绿)\n"
     "\n"
     "参数/术语:\n"
     "  - lag-1 自相关 = 相邻步的相关性\n"
     "  - lag-30 = 30 步外的相关性 (影响 horizon 上限)\n"
     "\n"
     "数据对比:\n"
     "  原始信号:  lag-1=0.87, LSTM=0.111, TFM=0.084\n"
     "  移动均值:  lag-1=1.00, LSTM=0.067, TFM=0.038\n"
     "  SavGol ★: lag-1=1.00, LSTM=0.055, TFM=0.044\n"
     "  Butter:   lag-1=1.00, LSTM=0.107, TFM=0.060"),

    ("图 4-3", "方案 B - 分解后预测曲线样例",
     "SavGol 趋势行: 预测曲线终于能跟上波动",
     "分解_3_预测曲线对比.png",
     "信号分解后预测视觉效果明显改善",
     "图表内容:\n"
     "  4 行 × 3 列, 每行一种信号, 每列一个样例\n"
     "  行1 原始 / 行2 移动均值 / 行3 SavGol / 行4 Butter\n"
     "\n"
     "观察:\n"
     "  - 原始信号行: 真实值高频抖动, 预测平均线\n"
     "  - SavGol 行 (样例 2/3): 预测曲线弯曲跟随 ✓\n"
     "  - horizon 仍是 32, 所以仍有改进空间"),

    ("图 5-1", "缩短预测窗口 horizon 的视觉效果",
     "horizon = 4 / 8 / 16 / 32 在 SavGol 趋势上",
     "horizon_1_缩短窗口对比.png",
     "horizon=16 是最佳平衡: 预测曲线明显跟随, MAE 低",
     "图表内容:\n"
     "  4 行 × 4 列, 每行一个 horizon, 每列一个样例\n"
     "  可见:\n"
     "   h=4 太短看不出\n"
     "   h=8 TimesFM 开始跟随\n"
     "   h=16 ★ 两个模型都完美跟随真实弧度\n"
     "   h=32 又退化为直线\n"
     "\n"
     "参数/术语:\n"
     "  - horizon: 预测未来多少秒\n"
     "  - context: 固定 256 秒历史\n"
     "  - 信号: SavGol(61,3) 分解的趋势部分"),

    ("图 5-2", "预测窗口与精度的定量关系",
     "MAE 随 horizon 指数增长, TimesFM 始终领先",
     "horizon_2_MAE曲线.png",
     "horizon = 16 相对误差仅 2.1%, 推荐工程使用",
     "图表内容:\n"
     "  MAE vs horizon 折线图\n"
     "  红色 = LSTM, 绿色 = TimesFM\n"
     "\n"
     "精度对照:\n"
     "  h=4  TFM MAE=0.005 (相对误差 0.5%)\n"
     "  h=8  TFM MAE=0.009 (相对误差 1.0%)\n"
     "  h=16 TFM MAE=0.018 (相对误差 2.1%) ← 推荐\n"
     "  h=32 TFM MAE=0.051 (相对误差 5.9%)\n"
     "\n"
     "参数/术语:\n"
     "  - 相对误差 = MAE / 信号均值 × 100%\n"
     "    < 5% 视为可接受, < 2% 视为优秀"),

    ("图 6-1", "多属性泛化性 — 预测曲线",
     "最佳配置 (SavGol + h=16) 在 4 个代表属性上",
     "多属性_1_预测曲线.png",
     "连续类 3 属性效果佳, 脉冲类 ShockZ 不适用",
     "图表内容:\n"
     "  4 行 × 3 列, 每行一个属性, 每列一个样例\n"
     "  属性顺序: Magnitude / aRMSX / vRMSM / ShockZ\n"
     "\n"
     "观察:\n"
     "  - Magnitude/aRMSX/vRMSM: 预测明显跟随真实值 ✓\n"
     "  - ShockZ: 预测偏平 (脉冲型不适合此方案)\n"
     "\n"
     "参数/术语:\n"
     "  - SavGol(61,3): 61秒窗口, 3阶多项式\n"
     "  - horizon = 16 秒\n"
     "  - context = 256 秒"),

    ("图 6-2", "多属性泛化性 — 指标汇总",
     "TimesFM 在连续类属性上相对误差稳定 2.1%",
     "多属性_2_指标对比.png",
     "TimesFM 跨属性一致性远胜 LSTM, 零样本优势显著",
     "图表内容:\n"
     "  左: MAE 对比 (log 轴)\n"
     "  中: 相对误差% 对比\n"
     "  右: 趋势信号 ACF 对比 (lag-1 蓝, lag-16 紫)\n"
     "\n"
     "数据表:\n"
     "  Magnitude: LSTM 3.8%, TFM 2.1%, ACF-16=0.81\n"
     "  aRMSX:    LSTM 4.0%, TFM 2.2%, ACF-16=0.81\n"
     "  vRMSM:    LSTM 4.2%, TFM 2.1%, ACF-16=0.82\n"
     "  ShockZ:   LSTM 2.8%, TFM 2.4%, ACF-16=0.56 ← 弱\n"
     "\n"
     "结论:\n"
     "  - TimesFM 跨属性相对误差极稳定 (2.1-2.4%)\n"
     "  - LSTM 差异较大 (2.8-4.2%)\n"
     "  - ShockZ ACF 只有 0.56, 预示 SavGol 对它失效"),

    ("总结", "项目结论与工程建议",
     "5 项关键发现 + 4 种业务场景的推荐配置",
     None,
     None,
     "发现 1: 1s 尺度原始振动接近高频噪声, 逐点预测必然退化\n"
     "发现 2: SavGol 信号分解能有效分离可预测趋势\n"
     "发现 3: horizon=16 秒是预测精度和业务价值的最佳平衡\n"
     "发现 4: TimesFM 在 6/7 参数组合和 3/4 属性上优于 LSTM\n"
     "发现 5: Shock 类脉冲信号不适合趋势预测, 需用事件检测\n"
     "\n"
     "工程建议:\n"
     "  场景1 振动趋势监测: SavGol(61,3) + TimesFM + h=16, 相对误差 2%\n"
     "  场景2 短期实时预测: SavGol(31,2) + TimesFM + h=8,  相对误差 1%\n"
     "  场景3 异常事件检测: 原始信号 + LSTM + 残差 3 级告警\n"
     "  场景4 冲击监控:     ShockX/Y/Z + 脉冲事件检测算法"),
]


# ============================================================
# 封面 / 目录 / 总结页 特殊渲染函数
# ============================================================
def render_cover(slide, title, subtitle, note_text):
    set_slide_bg(slide, WHITE)
    # 彩色斜条装饰
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.25), fill=PRIMARY)
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), fill=ACCENT)

    add_text(slide, Inches(0.5), Inches(2.3), SLIDE_W - Inches(1), Inches(1.4),
             title, size=40, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    # 装饰横线
    add_rect(slide, Inches(5), Inches(3.9), Inches(3.3), Inches(0.04),
             fill=ACCENT)

    add_text(slide, Inches(0.5), Inches(4.1), SLIDE_W - Inches(1), Inches(0.5),
             subtitle, size=18, color=TEXT_MID, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(5.3), SLIDE_W - Inches(1), Inches(0.4),
             "数据: MPB01_6.11 08_18.csv  ·  设备: MPB_01  ·  11 属性时序",
             size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)

    add_notes(slide, note_text)


def render_agenda(slide, title, subtitle, note_text):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.5), SLIDE_W - Inches(1), Inches(0.7),
             title, size=32, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(1.3), SLIDE_W - Inches(1), Inches(0.4),
             subtitle, size=14, color=TEXT_MID, align=PP_ALIGN.CENTER)

    chapters = [
        ("01", "数据概览与可视化", "11 属性统计与时序分析", PRIMARY),
        ("02", "两种算法多参数对比", "LSTM vs TimesFM 网格实验", ACCENT),
        ("03", "预测直线问题诊断", "ACF 与朴素基线分析", SUCCESS),
        ("04", "信号分解与窗口优化", "SavGol 滤波 + 短 horizon 方案", PRIMARY),
        ("05", "多属性泛化与结论", "最佳配置的普适验证", ACCENT),
    ]
    y = Inches(2.3)
    for num, ctitle, csub, color in chapters:
        add_rect(slide, Inches(1.5), y, Inches(1.0), Inches(0.8),
                 fill=color)
        add_text(slide, Inches(1.5), y, Inches(1.0), Inches(0.8),
                 num, size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(2.5), y, Inches(9.5), Inches(0.8),
                 fill=LIGHT)
        add_text(slide, Inches(2.7), y + Inches(0.1),
                 Inches(9.3), Inches(0.4),
                 ctitle, size=16, bold=True, color=PRIMARY)
        add_text(slide, Inches(2.7), y + Inches(0.48),
                 Inches(9.3), Inches(0.3),
                 csub, size=11, color=TEXT_MID)
        y = y + Inches(0.95)

    add_notes(slide, note_text)


def render_summary(slide, title, subtitle, note_text):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.35), SLIDE_W - Inches(1), Inches(0.6),
             title, size=28, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.95), SLIDE_W - Inches(1), Inches(0.3),
             subtitle, size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)

    # 5 项关键发现
    findings = [
        "1s 尺度原始振动信号接近高频噪声, 逐点预测必然退化",
        "SavGol 信号分解能有效分离可预测的「趋势」部分",
        "horizon = 16 秒是预测精度与业务价值的最佳平衡",
        "TimesFM 在大多数配置和属性上稳定胜出 LSTM",
        "Shock 类脉冲信号不适合趋势预测, 需用事件检测方法",
    ]

    y = Inches(1.5)
    for i, f in enumerate(findings):
        add_rect(slide, Inches(0.5), y, Inches(1.2), Inches(0.5),
                 fill=PRIMARY if i % 2 == 0 else ACCENT)
        add_text(slide, Inches(0.5), y, Inches(1.2), Inches(0.5),
                 f"发现{i+1}", size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(1.7), y, Inches(11.0), Inches(0.5),
                 fill=LIGHT)
        add_text(slide, Inches(1.9), y + Inches(0.05),
                 Inches(10.8), Inches(0.4),
                 f, size=13, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.6)

    # 建议栏
    add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.3),
             fill=LIGHT)
    add_text(slide, Inches(0.8), Inches(4.85),
             Inches(12), Inches(0.4),
             "▍工程部署推荐配置",
             size=16, bold=True, color=PRIMARY)

    rec_lines = [
        "场景 1 · 振动趋势监测:  SavGol(61,3) + TimesFM + h=16, 相对误差 ~2%",
        "场景 2 · 短期实时预测:  SavGol(31,2) + TimesFM + h=8,  相对误差 ~1%",
        "场景 3 · 异常事件检测:  原始信号 + LSTM + 残差 3 级告警",
        "场景 4 · 冲击信号监控:  ShockX/Y/Z + 脉冲事件检测 (非时序预测)",
    ]
    y2 = Inches(5.3)
    for line in rec_lines:
        add_text(slide, Inches(1.0), y2, Inches(11.5), Inches(0.4),
                 "▶ " + line, size=12, color=TEXT_DARK)
        y2 = y2 + Inches(0.4)

    add_notes(slide, note_text)


# ============================================================
# 主渲染
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    total = len(SLIDES_SPEC)
    for page_i, spec in enumerate(SLIDES_SPEC):
        chart_no, title, subtitle, img, key_point, notes = spec
        s = prs.slides.add_slide(BLANK)

        if chart_no == "COVER":
            render_cover(s, title, subtitle, notes)
            continue
        if chart_no == "AGENDA":
            render_agenda(s, title, subtitle, notes)
            continue
        if chart_no == "总结":
            render_summary(s, title, subtitle, notes)
            continue

        # 普通图表页
        set_slide_bg(s, WHITE)
        slide_header(s, chart_no, title, subtitle)
        if img:
            add_figure_area(s, os.path.join(OUTPUT_DIR, img),
                            top=Inches(1.35), max_height=Inches(5.0))
        if key_point:
            add_key_point(s, key_point, top=Inches(6.5))
        add_footer(s, page_i + 1, total)
        add_notes(s, notes)

    prs.save(PPT_PATH)
    print(f"[OK] 已生成: {PPT_PATH}")
    print(f"总页数: {len(prs.slides)}")
    print(f"每页都含演讲者备注 (参数/术语/解读放在备注里)")


if __name__ == "__main__":
    build()
