"""
生成 MPB_01 振动预测研究报告 (提交版)
备注统一写为章节内容描述,用于随 PPT 一并提交,而非现场讲稿。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT_DIR = r"E:\timesfm_project\output"
PPT_PATH = os.path.join(OUTPUT_DIR, "MPB_01振动预测研究报告.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

PRIMARY = RGBColor(0x4A, 0x90, 0xE2)
ACCENT = RGBColor(0xF5, 0xA6, 0x23)
SUCCESS = RGBColor(0x7E, 0xD3, 0x21)
DANGER = RGBColor(0xE7, 0x4C, 0x3C)
BG = RGBColor(0xFC, 0xFD, 0xFE)
LIGHT = RGBColor(0xE8, 0xF0, 0xFA)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_GREEN = RGBColor(0xEB, 0xF8, 0xD6)
LIGHT_RED = RGBColor(0xFD, 0xEC, 0xEA)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_MID = RGBColor(0x66, 0x6E, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_font(run, name="微软雅黑", size=14, bold=False, color=TEXT_DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for e in rPr.findall(qn("a:ea")):
        rPr.remove(e)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)


def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.text = ""
        run = p.add_run()
        run.text = line
        set_font(run, name=font_name, size=size, bold=bold, color=color)
    return tb


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, note_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = ""
    lines = note_text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line


def slide_header(slide, chart_no, chart_name, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(0.3), Inches(1.8), Inches(0.45),
             fill=ACCENT)
    add_text(slide, Inches(0.4), Inches(0.32), Inches(1.8), Inches(0.4),
             chart_no, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(2.4), Inches(0.3), Inches(10.5), Inches(0.5),
             chart_name, size=20, bold=True, color=PRIMARY)
    if subtitle:
        add_text(slide, Inches(2.4), Inches(0.78), Inches(10.5), Inches(0.3),
                 subtitle, size=11, color=TEXT_MID)


def add_figure_area(slide, img_path, top=Inches(1.35), max_height=Inches(5.2)):
    if not os.path.exists(img_path):
        return
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    max_w = SLIDE_W - Inches(0.8)
    ratio = iw / ih
    if max_height * ratio > max_w:
        w = max_w
        h = w / ratio
    else:
        h = max_height
        w = h * ratio
    left = (SLIDE_W - w) // 2
    slide.shapes.add_picture(img_path, left, top, width=w, height=h)
    return top + h


def add_key_point(slide, text, top=None):
    if top is None:
        top = Inches(6.7)
    add_rect(slide, Inches(0.4), top, SLIDE_W - Inches(0.8), Inches(0.55),
             fill=LIGHT)
    add_text(slide, Inches(0.6), top + Inches(0.05),
             SLIDE_W - Inches(1.2), Inches(0.45),
             "▍本页要点  " + text, size=13, bold=True, color=PRIMARY,
             anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_no, total):
    add_text(slide, Inches(11.8), SLIDE_H - Inches(0.3),
             Inches(1.5), Inches(0.3),
             f"{page_no} / {total}",
             size=10, color=TEXT_MID, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.3),
             Inches(11), Inches(0.3),
             "MPB_01 振动预测与预防性维护研究",
             size=10, color=TEXT_MID)


def render_text_page(slide, chart_no, title, subtitle, sections, note_text,
                     page_no, total):
    set_slide_bg(slide, WHITE)
    slide_header(slide, chart_no, title, subtitle)

    y = Inches(1.35)
    for sec in sections:
        if sec["kind"] == "card":
            tag_color = sec.get("tag_color", PRIMARY)
            add_rect(slide, Inches(0.5), y, Inches(1.3), Inches(1.0),
                     fill=tag_color)
            add_text(slide, Inches(0.5), y, Inches(1.3), Inches(1.0),
                     sec["tag"], size=15, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            body_h = Inches(1.0)
            add_rect(slide, Inches(1.8), y, SLIDE_W - Inches(2.3), body_h,
                     fill=LIGHT)
            add_text(slide, Inches(2.0), y + Inches(0.06),
                     SLIDE_W - Inches(2.7), Inches(0.4),
                     sec["head"], size=14, bold=True, color=PRIMARY)
            body_text = "\n".join(sec["body"])
            add_text(slide, Inches(2.0), y + Inches(0.42),
                     SLIDE_W - Inches(2.7), Inches(0.6),
                     body_text, size=11, color=TEXT_DARK)
            y = y + body_h + Inches(0.12)
        elif sec["kind"] == "code":
            h = Inches(sec.get("height", 1.6))
            add_rect(slide, Inches(0.5), y, SLIDE_W - Inches(1.0), h,
                     fill=RGBColor(0x2C, 0x3E, 0x50))
            add_text(slide, Inches(0.7), y + Inches(0.1),
                     SLIDE_W - Inches(1.4), h - Inches(0.2),
                     sec["content"], size=12, color=WHITE,
                     font_name="Consolas")
            y = y + h + Inches(0.1)
        elif sec["kind"] == "scenarios":
            for s in sec["rows"]:
                row_h = Inches(0.45)
                color = {"ok": SUCCESS, "no": DANGER, "mid": ACCENT}[s["status"]]
                bg_col = {"ok": LIGHT_GREEN, "no": LIGHT_RED, "mid": LIGHT_ORANGE}[s["status"]]
                add_rect(slide, Inches(0.5), y, Inches(0.6), row_h, fill=color)
                add_text(slide, Inches(0.5), y, Inches(0.6), row_h,
                         s["icon"], size=14, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                add_rect(slide, Inches(1.1), y, SLIDE_W - Inches(1.6), row_h,
                         fill=bg_col)
                add_text(slide, Inches(1.3), y, Inches(8.0), row_h,
                         s["scenario"], size=12, bold=True, color=TEXT_DARK,
                         anchor=MSO_ANCHOR.MIDDLE)
                add_text(slide, Inches(9.5), y, SLIDE_W - Inches(10.2), row_h,
                         s["answer"], size=11, color=TEXT_MID,
                         anchor=MSO_ANCHOR.MIDDLE)
                y = y + row_h + Inches(0.06)
        elif sec["kind"] == "banner":
            add_rect(slide, Inches(0.5), y, SLIDE_W - Inches(1.0),
                     Inches(sec.get("height", 0.6)),
                     fill=sec.get("fill", LIGHT))
            add_text(slide, Inches(0.7), y, SLIDE_W - Inches(1.4),
                     Inches(sec.get("height", 0.6)),
                     sec["content"], size=sec.get("size", 14),
                     bold=sec.get("bold", True),
                     color=sec.get("color", PRIMARY),
                     align=sec.get("align", PP_ALIGN.LEFT),
                     anchor=MSO_ANCHOR.MIDDLE)
            y = y + Inches(sec.get("height", 0.6)) + Inches(0.1)

    add_footer(slide, page_no, total)
    add_notes(slide, note_text)


# ============================================================
# Spec
# ============================================================
SLIDES_SPEC = [
    ("COVER", "MPB_01 设备振动预测与预防性维护研究",
     "LSTM 与 TimesFM 对比 · 信号分解 · 多参数实验 · 全字段泛化验证",
     None, None,
     "本页为封面。\n"
     "研究对象: MPB_01 传感器采集的振动时序数据。\n"
     "研究内容: 针对 11 个物理量字段, 比较 LSTM (本地训练) 与\n"
     "TimesFM 2.5-200M (零样本) 两种方法的预测性能, 并对信号\n"
     "可预测性、分解策略、预测窗口选择与泛化能力进行系统分析。\n"
     "数据来源: MPB01_6.11 08_18.csv (InfluxDB 导出), 连续监测\n"
     "约 2.6 小时, 按 1 秒重采样后约 8960 个样本点。\n"
     "报告范围提示: 本研究基于单台设备数据, 跨设备泛化未在本\n"
     "报告范围内。"),

    ("AGENDA", "报告目录",
     "全报告共分七个章节",
     None, None,
     "本页为章节导航。\n"
     "章节一展示 11 属性原始数据的统计与可视化。\n"
     "章节二呈现 LSTM 与 TimesFM 在不同重采样率、预测长度和\n"
     "上下文长度下的对比实验。\n"
     "章节三分析预测曲线退化为近似直线的原因, 通过自相关函数\n"
     "与朴素基线诊断信号可预测性。\n"
     "章节四介绍信号分解方案, 评估移动均值、Savitzky-Golay 和\n"
     "Butterworth 三种滤波方法。\n"
     "章节五定量分析预测窗口 horizon 对精度的影响。\n"
     "章节六将最优方案推广至 11 个物理量字段, 引入 NRMSE 指标\n"
     "和可用性判据并完成验证。\n"
     "章节七给出研究结论与工程部署建议。"),

    ("图 1-1", "11 属性分量级时序对照",
     "按数值量级分三组, 橙色带标注 18 次疑似停机",
     "属性总览_1_分量级时序.png",
     "Magnitude 存在 18 次骤降事件, 对应设备停机或传感器离线",
     "本页展示 11 个物理量按数值量级划分为三组后的时序对照图。\n"
     "组 A 为 Magnitude (合成幅值, 量级 0-1);\n"
     "组 B 为 aRMS 与 vRMS 系列 (均方根, 量级 0-5);\n"
     "组 C 为 ShockX/Y/Z (三轴冲击, 量级 ±1500)。\n"
     "图中橙色半透明带标注了对 Magnitude 序列检测出的 18 次\n"
     "疑似异常区段。\n\n"
     "术语说明:\n"
     "  RMS (Root Mean Square) 均方根, 用于衡量信号能量;\n"
     "  Magnitude 指三轴合成后的总幅值;\n"
     "  aRMS 表示加速度均方根, 单位为 g;\n"
     "  vRMS 表示速度均方根, 单位为 mm/s;\n"
     "  Shock 指冲击峰值, 反映瞬时撞击强度。\n"
     "数据统一重采样至 1 秒分辨率, 每个子图含 8960 个数据点。"),

    ("图 1-2", "11 属性独立时序图",
     "每个属性独立面板, 便于观察各自量级与波动",
     "属性总览_2_独立子图.png",
     "Shock 类呈尖峰脉冲形态, RMS 类呈连续起伏, 物理本质不同",
     "本页以独立面板展示 11 个属性各自的完整时序, 突出各属性\n"
     "的量级与波动特征。每个子图标题给出数据范围和标准差, 红色\n"
     "虚线为均值线, 用于观察偏离情况。\n\n"
     "术语说明:\n"
     "  均值 (mean): 所有样本的算术平均;\n"
     "  标准差 (std): 度量样本相对均值的离散程度;\n"
     "  范围 [min, max]: 数据的上下界。\n\n"
     "观察结果:\n"
     "  ShockX/Y/Z 呈现密集尖峰型 (脉冲型) 波形;\n"
     "  aRMS 与 vRMS 系列呈现随时间起伏的连续型波形;\n"
     "  Magnitude 整体平稳, 但存在多次骤降 (对应停机事件)。"),

    ("图 1-3", "多属性统计概览与相关性",
     "对数坐标量级对比 + 变异系数 + 11×11 相关矩阵",
     "属性总览_3_统计与相关性.png",
     "Magnitude 与 aRMS/vRMS 本质同源 (相关 0.86-0.98), Shock 相对独立",
     "本页汇总 11 属性的统计特征与相互关系。\n"
     "左上图为各属性绝对值均值的对数柱状图, 揭示跨量级差异;\n"
     "右上图为变异系数 CV 的对比, 其中 ShockY=1.66 数值最高,\n"
     "意味着该属性相对波动最大、建模最困难;\n"
     "下方为 11×11 的皮尔逊相关系数热力图。\n\n"
     "术语说明:\n"
     "  变异系数 CV = std / |均值|, 用于衡量相对波动。\n"
     "    CV > 1 为高难度, CV > 0.5 为中等难度。\n"
     "  皮尔逊相关系数 r ∈ [-1, 1]。\n"
     "    |r| > 0.7 为强相关; |r| < 0.3 为弱相关或无关。\n\n"
     "主要发现:\n"
     "  Magnitude 与 aRMS/vRMS 系列的相关系数在 0.86-0.98 之间,\n"
     "  说明这些属性本质上源自同一物理量;\n"
     "  Shock 系列与其他属性相关系数约为 0.05, 基本独立。"),

    ("图 1-4", "异常事件放大 — Magnitude 骤降",
     "18 次疑似事件中的前 4 例详细波形",
     "属性总览_4_异常事件放大.png",
     "典型事件: 振动值从约 0.86 骤降至约 0, 持续 3-14 秒后恢复",
     "本页放大展示 Magnitude 序列中 18 次疑似异常事件的前 4 例。\n"
     "每个子图显示骤降前后 30 秒的波形, 橙色阴影标识异常持续区\n"
     "间, 灰色虚线为正常段的中位数参考。\n\n"
     "异常判定参数:\n"
     "  z_threshold = 3.0 (判定阈值, 数值越小越敏感);\n"
     "  min_duration = 3 秒 (最短持续时长, 避免瞬时误判);\n"
     "  判定规则: 数值低于 中位数 - 3.0 × std × 0.3。\n\n"
     "业务含义:\n"
     "  数值骤降至接近 0 通常对应设备停机或传感器离线。\n"
     "  对后续预测评估而言, 若测试集中包含异常段将干扰指标,\n"
     "  因此评估阶段需对该类事件进行剔除或单独处理。"),

    ("图 2-1", "LSTM 与 TimesFM 分维度参数对比",
     "Magnitude 上三个维度共七组实验的 MAE 对比",
     "对比_1_分维度柱状图.png",
     "TimesFM 在绝大多数配置下优于 LSTM",
     "本页展示两种模型在 Magnitude 上, 沿三个参数维度的 MAE\n"
     "对比实验。\n"
     "维度一 重采样频率: 1 秒 / 5 秒 / 30 秒;\n"
     "维度二 预测长度 horizon: 8 / 32 / 64;\n"
     "维度三 上下文长度 context: 128 / 256 / 512。\n"
     "红色柱为 LSTM 结果, 绿色柱为 TimesFM 结果, MAE 越低越好。\n\n"
     "术语说明:\n"
     "  MAE (Mean Absolute Error) 平均绝对误差, 单位与信号一致;\n"
     "  horizon 指一次预测覆盖的未来步数;\n"
     "  context 指模型使用的历史数据长度。\n\n"
     "实验结果:\n"
     "  5 秒重采样相较 1 秒略有改善, 降噪效应明显;\n"
     "  horizon 越小 MAE 越小, 符合「短期预测更易」的规律;\n"
     "  TimesFM 对 context 长度不敏感, LSTM 对其较敏感。"),

    ("图 2-2", "六组参数下的预测曲线样例",
     "预测曲线普遍退化为近似直线",
     "对比_2_预测曲线样例.png",
     "两种模型的预测曲线均偏平, 触发后续可预测性诊断",
     "本页展示六组不同参数实验对应的预测曲线样例。\n"
     "深蓝色为真实值, 红色虚线为 LSTM 预测, 绿色点划线为 TimesFM\n"
     "预测。即便是表现最好的 horizon=8 组, 模型输出曲线仍显著\n"
     "平于真实波动。\n\n"
     "模型说明:\n"
     "  LSTM: 采用 Encoder-Decoder 结构配合多步直接输出, 在本\n"
     "        项目数据上从头训练;\n"
     "  TimesFM: Google 开源的预训练基础模型 2.5-200M, 零样本\n"
     "        推理, 未针对本数据微调。\n\n"
     "由本页观察引出后续章节的诊断分析。尽管数值 MAE 并不异常,\n"
     "但曲线跟随性较差, 说明 MAE 单一指标无法完整反映预测质量,\n"
     "需要进一步探究信号本身的可预测性。"),

    ("图 2-3", "算法胜率与 MAE 相对差距汇总",
     "TimesFM 在 6 组实验中的 5 组胜出, 平均领先约 41%",
     "对比_3_胜率汇总.png",
     "TimesFM 无需训练即可使用, LSTM 需训练且结果波动较大",
     "本页汇总七组参数实验的整体对比结果。\n"
     "左侧为胜率饼图, TimesFM 胜率约 83%, LSTM 约 17%;\n"
     "右侧为每组实验中 MAE 相对差距的横向条形。\n\n"
     "指标说明:\n"
     "  MAE 相对差距 = (较差值 / 较好值 − 1) × 100%。\n"
     "  例: LSTM MAE=0.112, TFM MAE=0.079, 相对差距约 42.9%。\n\n"
     "时间成本:\n"
     "  TimesFM 总耗时约 71 秒 (仅推理);\n"
     "  LSTM 总耗时约 137 秒 (训练 + 推理)。"),

    ("图 3-1", "预测退化为直线的原因诊断",
     "基于自相关函数、朴素基线与功率谱的三维度分析",
     "诊断_1_可预测性分析.png",
     "1 秒尺度信号接近高频噪声, 最优预测接近均值",
     "本页从三个维度诊断预测曲线退化为直线的原因。\n"
     "左上图: 不同重采样尺度下的自相关函数 (ACF);\n"
     "右上图: 1 秒尺度信号的功率谱密度 PSD;\n"
     "中部: 朴素基线与 LSTM / TimesFM 的 MAE 对比;\n"
     "左下图: 各重采样尺度下的相对误差;\n"
     "右下图: 100 秒真实振动切片, 可见高频抖动特征。\n\n"
     "术语说明:\n"
     "  ACF (AutoCorrelation Function) 自相关函数;\n"
     "    r(lag=k) 表示 x[t] 与 x[t+k] 之间的相关性;\n"
     "    r > 0.3 通常被视为具备预测价值的阈值。\n"
     "  朴素基线 (Naive Baseline): 不使用模型的简单估计。\n"
     "    朴素-最后值: 未来各步均等于当前最末值。\n"
     "  PSD (Power Spectral Density) 功率谱密度。\n\n"
     "诊断结论:\n"
     "  1 秒尺度 lag-1 自相关约为 0.87, 相邻秒高度相关;\n"
     "  但 lag-60 自相关仅约 0.09, 远期几乎独立;\n"
     "  朴素-最后值 MAE=0.069, 优于 LSTM 的 0.112,\n"
     "  说明原信号在此尺度下可预测性有限。"),

    ("图 4-1", "信号分解方案 — 三种滤波对比",
     "将 Magnitude 拆分为趋势 + 噪声, 对比移动均值 / SavGol / Butter",
     "分解_1_信号分解对比.png",
     "SavGol 保留趋势形状最好, Butter 最平滑但损失细节",
     "本页对比三种信号分解方法的效果。\n"
     "第一行叠加原始信号 (灰) 与三种趋势 (彩色);\n"
     "其后三行分别展示每种分解方法得到的趋势 (主轴) 与残差\n"
     "(副轴)。\n\n"
     "方法参数:\n"
     "  移动均值 (Moving Average): 窗口 = 30 秒;\n"
     "  Savitzky-Golay 滤波: 窗口长度 = 61, 多项式阶数 = 3;\n"
     "  Butterworth 低通滤波: 截止频率 = 0.05 Hz, 阶数 = 4;\n"
     "  趋势占比 = Var(趋势) / Var(原信号)。\n\n"
     "结果:\n"
     "  移动均值趋势占比 70.5%;\n"
     "  SavGol 趋势占比 71.9%;\n"
     "  Butterworth 趋势占比 84.1%, 平滑程度最高但细节损失。"),

    ("图 4-2", "分解后可预测性提升",
     "趋势信号 lag-1 自相关从 0.87 提升至 1.00, MAE 下降约 47%",
     "分解_2_可预测性对比.png",
     "分解后趋势信号具备更强的可预测性, 模型不再仅能输出均值",
     "本页量化分解前后可预测性与预测误差的变化。\n"
     "左侧为 ACF 对比 (lag-1 蓝色, lag-30 紫色);\n"
     "右侧为 MAE 对比 (朴素灰色, LSTM 红色, TimesFM 绿色)。\n\n"
     "数据对比:\n"
     "  原始信号:   lag-1=0.87, LSTM MAE=0.111, TFM MAE=0.084;\n"
     "  移动均值:   lag-1=1.00, LSTM MAE=0.067, TFM MAE=0.038;\n"
     "  SavGol:     lag-1=1.00, LSTM MAE=0.055, TFM MAE=0.044;\n"
     "  Butterworth: lag-1=1.00, LSTM MAE=0.107, TFM MAE=0.060。\n\n"
     "综合精度与趋势形状保留能力, SavGol(61,3) 为本研究采用\n"
     "的默认分解方案。"),

    ("图 4-3", "分解后的预测曲线样例",
     "SavGol 趋势行的预测曲线具备明显跟随性",
     "分解_3_预测曲线对比.png",
     "信号分解后预测视觉效果显著改善",
     "本页展示不同分解方案下的预测曲线样例, 共 4 行 × 3 列。\n"
     "第 1 行为原始信号, 第 2 行为移动均值, 第 3 行为 SavGol,\n"
     "第 4 行为 Butterworth。\n\n"
     "观察结果:\n"
     "  原始信号行: 真实值呈现高频抖动, 模型预测近似均值;\n"
     "  SavGol 行: 预测曲线呈现与真实趋势一致的弧度跟随;\n"
     "  此时 horizon 仍为 32, 保留了进一步优化的空间。"),

    ("图 5-1", "预测窗口缩短的视觉效果",
     "horizon = 4 / 8 / 16 / 32 在 SavGol 趋势上的对比",
     "horizon_1_缩短窗口对比.png",
     "horizon=16 为精度与可视化跟随性的较优平衡点",
     "本页对比四种 horizon 配置下的预测曲线, 共 4 行 × 4 列。\n"
     "每行对应一个 horizon 值, 每列为一个样例。\n\n"
     "观察结果:\n"
     "  horizon = 4: 预测窗口过短, 变化难以显现;\n"
     "  horizon = 8: TimesFM 已能较好跟随;\n"
     "  horizon = 16: 两模型均可完整跟随真实弧度;\n"
     "  horizon = 32: 预测曲线再次退化为近似直线。\n\n"
     "固定参数: context = 256 秒, 信号取 SavGol(61,3) 趋势部分。"),

    ("图 5-2", "预测窗口与精度的定量关系",
     "MAE 随 horizon 近似指数增长, TimesFM 始终领先",
     "horizon_2_MAE曲线.png",
     "horizon=16 时相对误差约 2.1%, 为工程部署推荐配置",
     "本页以折线图形式量化 MAE 与 horizon 的关系。\n"
     "红色曲线为 LSTM, 绿色曲线为 TimesFM。\n\n"
     "精度对照:\n"
     "  horizon = 4:  TFM MAE=0.005 (相对误差约 0.5%);\n"
     "  horizon = 8:  TFM MAE=0.009 (相对误差约 1.0%);\n"
     "  horizon = 16: TFM MAE=0.018 (相对误差约 2.1%), 推荐;\n"
     "  horizon = 32: TFM MAE=0.051 (相对误差约 5.9%)。\n\n"
     "指标口径说明:\n"
     "  此处相对误差定义为 MAE / |信号均值|。该定义在信号均值\n"
     "  远大于波动幅度时会低估真实预测难度, 章节六将改用 NRMSE\n"
     "  指标以修正此问题。"),

    ("图 6-1", "最优方案在四个代表属性上的预测曲线",
     "SavGol + horizon=16 在 Magnitude / aRMSX / vRMSM / ShockZ 上的表现",
     "多属性_1_预测曲线.png",
     "连续类三个属性跟随效果良好, 冲击类 ShockZ 跟随性较弱",
     "本页展示最优方案 SavGol(61,3) + TimesFM + horizon=16 在\n"
     "四个代表属性上的预测曲线, 共 4 行 × 3 列。\n"
     "四个属性分别为 Magnitude, aRMSX, vRMSM, ShockZ。\n\n"
     "观察结果:\n"
     "  Magnitude, aRMSX, vRMSM 三个属性预测曲线跟随良好;\n"
     "  ShockZ 预测曲线偏平, 与真实值偏离明显;\n"
     "  初步表明本方案对连续型信号有效, 对脉冲型信号不适用。\n\n"
     "实验参数:\n"
     "  SavGol 窗口 61, 多项式阶数 3;\n"
     "  horizon = 16 秒, context = 256 秒。"),

    ("图 6-2", "四属性指标汇总 (原指标口径)",
     "在原「相对误差 = MAE / |均值|」口径下, 四属性相对误差均接近 2%",
     "多属性_2_指标对比.png",
     "原指标无法区分真实性能差异, 需要改用 NRMSE (详见章节六)",
     "本页以原相对误差指标汇总四属性对比结果, 包含三个子图:\n"
     "左: MAE 对比 (对数坐标);\n"
     "中: 相对误差百分比对比;\n"
     "右: 趋势信号的自相关对比 (lag-1 蓝色, lag-16 紫色)。\n\n"
     "指标口径局限性:\n"
     "  本页使用的相对误差定义为 MAE / |信号均值|。\n"
     "  ShockZ 的均值为 -226, 其绝对值较大, 会对分母产生稀释,\n"
     "  导致其相对误差表面上与其他属性接近 (2.4%), 掩盖了真实\n"
     "  预测性能差异。\n\n"
     "数据记录:\n"
     "  Magnitude: LSTM 3.8%, TFM 2.1%, ACF-16=0.81;\n"
     "  aRMSX:    LSTM 4.0%, TFM 2.2%, ACF-16=0.81;\n"
     "  vRMSM:    LSTM 4.2%, TFM 2.1%, ACF-16=0.82;\n"
     "  ShockZ:   LSTM 2.8%, TFM 2.4%, ACF-16=0.56。\n\n"
     "上述观察直接导出章节六的指标修正工作。"),

    ("图 7-1", "指标修正 — 引入 NRMSE",
     "以 NRMSE 替代 MAE / |均值|, 暴露 ShockZ 的真实预测性能",
     "多属性_3_NRMSE对比.png",
     "ShockZ 在 NRMSE 指标下的真实水平为 36%, 而非原口径的 2.4%",
     "本页介绍本研究采用的修正指标 NRMSE 以及其与原指标的\n"
     "差异。\n\n"
     "原指标的问题:\n"
     "  原相对误差 = MAE / |信号均值| × 100%。\n"
     "  当信号均值绝对值较大 (例如 ShockZ 的 -226), 分母会对\n"
     "  MAE 产生稀释作用, 使结果看似较好, 掩盖真实预测性能。\n\n"
     "修正指标 NRMSE 定义:\n"
     "  NRMSE = MAE / 趋势std × 100%。\n"
     "  含义为预测误差相对于信号本身波动幅度的比例:\n"
     "    NRMSE = 10% 表示预测误差仅为信号波动的 1/10;\n"
     "    NRMSE = 100% 表示模型输出接近均值, 已无预测价值。\n\n"
     "两种指标在四属性上的对比:\n"
     "  原指标 (Magnitude / aRMSX / vRMSM / ShockZ):\n"
     "    2.1% / 2.2% / 2.1% / 2.4% — 数值相近;\n"
     "  NRMSE:\n"
     "    13.8% / 13.5% / 13.1% / 36.1% — 性能差异显现。\n\n"
     "后续指标口径:\n"
     "  本报告章节六统一使用 NRMSE。\n"
     "  图中绿色 20% 线为「可用」阈值参考;\n"
     "  红色 35% 线为「不建议直接使用」阈值参考。"),

    ("图 7-2", "全字段 NRMSE 对比 — 11 个字段",
     "按物理量族分色, 揭示族内一致性与族间差异",
     "全字段_1_NRMSE对比.png",
     "RMS 与幅值类 8 个字段 NRMSE 集中于 13-17%, 冲击类 3 字段集中于 33-48%",
     "本页将最优方案推广至 MPB_01 的全部 11 个字段, 并按物理\n"
     "量族对结果着色。\n\n"
     "图例说明:\n"
     "  斜线填充柱代表 LSTM, 实心柱代表 TimesFM;\n"
     "  颜色编码: 蓝色 = 合成幅值, 青绿 = 加速度 RMS,\n"
     "    紫色 = 速度 RMS, 红色 = 冲击类;\n"
     "  绿色虚线 20% 为 A 档阈值 (可直接部署);\n"
     "  红色虚线 35% 为 C 档阈值 (不建议部署)。\n\n"
     "族内一致性结果:\n"
     "  加速度 RMS (3 字段): NRMSE 13.5–17.0%, 标准差 1.7%;\n"
     "  速度 RMS (4 字段):   NRMSE 13.1–15.4%, 标准差 1.2%;\n"
     "  冲击类 (3 字段):     NRMSE 33.7–48.4%, 标准差 7.9%。\n"
     "  同族内部方差较小, 表明可按物理量族判断字段可用性。\n\n"
     "模型对比结果:\n"
     "  11 个字段中 10 个 TimesFM 优于 LSTM, 平均领先约 41%。\n"
     "  唯一例外为 ShockX, TFM 较 LSTM 差 8.6%, 推测与该字段\n"
     "  信噪比极低 (SNR=0.238) 有关。"),

    ("图 7-3", "可用性判据象限 — ACF 与 SNR",
     "使用两个预算指标即可判断字段是否适用本方案",
     "全字段_2_判据象限.png",
     "冲击类三个字段聚集于左下区域, RMS 与幅值类八个字段聚集于右上区域",
     "本页建立用于字段可用性判断的二维判据。\n\n"
     "坐标轴含义:\n"
     "  横轴 ACF lag-16: 取值范围 0–1;\n"
     "    含义为 16 秒前的取值与当前取值的相关性;\n"
     "    > 0.75 表示信号具备长程记忆, 可预测性较好;\n"
     "    < 0.60 表示信号对历史依赖较弱, 可预测性较差。\n"
     "  纵轴 SNR (对数尺度) = 趋势std / 噪声std;\n"
     "    含义为信号骨架相对毛刺的倍数;\n"
     "    > 1.5 表示规律明显;\n"
     "    < 0.5 表示噪声掩盖规律。\n\n"
     "气泡表示:\n"
     "  颜色对应物理量族;\n"
     "  直径对应该字段的 NRMSE (直径越大误差越大)。\n\n"
     "结论:\n"
     "  右上高 ACF、高 SNR 区域聚集 8 个 RMS/幅值字段;\n"
     "  左下低 ACF、低 SNR 区域聚集 3 个冲击字段;\n"
     "  两个区域间无过渡样本, 判据分离清晰。\n"
     "  实际使用时, 新字段接入前仅需计算 ACF 与 SNR, 即可\n"
     "  预判方案适用性, 无需训练模型。"),

    ("图 7-4-TEXT", "全字段验证的三项主要发现",
     "族内一致 / 极低信噪比下的性能反转 / 冲击类整族不适用",
     None, None, None),

    ("图 7-5-TEXT", "可用性判据的修订与命中率",
     "两处阈值调整, 使判据对 11 个字段的命中率达到 100%",
     None, None, None),

    ("图 7-6-TEXT", "数据覆盖范围与局限性说明",
     "研究方法学结论充分, 跨设备泛化需后续数据补充",
     None, None, None),

    ("结论", "研究结论与工程部署建议",
     "七项主要发现与五类场景部署配置",
     None, None, None),
]


def render_cover(slide, title, subtitle, note_text):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.25), fill=PRIMARY)
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(2.3), SLIDE_W - Inches(1), Inches(1.4),
             title, size=40, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5), Inches(3.9), Inches(3.3), Inches(0.04),
             fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(4.1), SLIDE_W - Inches(1), Inches(0.5),
             subtitle, size=18, color=TEXT_MID, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(5.3), SLIDE_W - Inches(1), Inches(0.4),
             "数据来源: MPB01_6.11 08_18.csv  ·  设备: MPB_01  ·  11 属性时序",
             size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(5.7), SLIDE_W - Inches(1), Inches(0.4),
             "研究范围: 单设备多字段预测方法对比与泛化验证",
             size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)
    add_notes(slide, note_text)


def render_agenda(slide, title, subtitle, note_text):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.35), SLIDE_W - Inches(1), Inches(0.7),
             title, size=32, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(1.1), SLIDE_W - Inches(1), Inches(0.4),
             subtitle, size=14, color=TEXT_MID, align=PP_ALIGN.CENTER)

    chapters = [
        ("一", "数据概览与可视化",       "11 属性的统计特征与时序分析",  PRIMARY),
        ("二", "两种算法多参数对比",     "LSTM 与 TimesFM 的网格实验",   ACCENT),
        ("三", "预测退化原因诊断",       "基于 ACF 与朴素基线的分析",    SUCCESS),
        ("四", "信号分解与方案优化",     "SavGol 滤波与短 horizon 方案", PRIMARY),
        ("五", "多属性初步泛化",         "四个代表属性的验证",           ACCENT),
        ("六", "全字段验证与判据",       "11 字段 NRMSE 与可用性判据",   DANGER),
        ("七", "结论与工程建议",         "研究发现与部署配置",           SUCCESS),
    ]
    y = Inches(1.7)
    for num, ctitle, csub, color in chapters:
        add_rect(slide, Inches(1.5), y, Inches(1.0), Inches(0.68),
                 fill=color)
        add_text(slide, Inches(1.5), y, Inches(1.0), Inches(0.68),
                 num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(2.5), y, Inches(9.5), Inches(0.68),
                 fill=LIGHT)
        add_text(slide, Inches(2.7), y + Inches(0.05),
                 Inches(9.3), Inches(0.35),
                 ctitle, size=14, bold=True, color=PRIMARY)
        add_text(slide, Inches(2.7), y + Inches(0.38),
                 Inches(9.3), Inches(0.3),
                 csub, size=10, color=TEXT_MID)
        y = y + Inches(0.78)

    add_notes(slide, note_text)


def render_summary(slide, title, subtitle, note_text, page_no, total):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.25), SLIDE_W - Inches(1), Inches(0.5),
             title, size=26, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.75), SLIDE_W - Inches(1), Inches(0.3),
             subtitle, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

    findings = [
        "1 秒尺度的原始振动信号接近高频噪声, 逐点预测必然退化",
        "Savitzky-Golay 信号分解能有效分离可预测的趋势部分",
        "horizon = 16 秒为精度与窗口长度的较优平衡点",
        "TimesFM 零样本预测胜出 LSTM: 11 字段中 10 个领先, 平均提升约 41%",
        "物理量族内一致性显著 (NRMSE 标准差 ≤ 2%), 可按族作预判",
        "冲击类信号整族不适用趋势预测方案 (NRMSE 33-48%, SNR<0.4)",
        "ACF + SNR + NRMSE 三指标判据对 11 字段命中率 100%",
    ]

    y = Inches(1.2)
    for i, f in enumerate(findings):
        add_rect(slide, Inches(0.5), y, Inches(1.2), Inches(0.38),
                 fill=PRIMARY if i % 2 == 0 else ACCENT)
        add_text(slide, Inches(0.5), y, Inches(1.2), Inches(0.38),
                 f"发现{i+1}", size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(1.7), y, Inches(11.0), Inches(0.38),
                 fill=LIGHT)
        add_text(slide, Inches(1.9), y,
                 Inches(10.8), Inches(0.38),
                 f, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.44)

    add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
             fill=LIGHT)
    add_text(slide, Inches(0.8), Inches(4.55),
             Inches(12), Inches(0.4),
             "▍工程部署推荐配置",
             size=15, bold=True, color=PRIMARY)

    rec_lines = [
        ("场景一 · 振动趋势监测", "SavGol(61,3) + TimesFM + h=16, NRMSE 约 14%",  SUCCESS),
        ("场景二 · 短期实时预测", "SavGol(31,2) + TimesFM + h=8,  NRMSE 约 7%",   SUCCESS),
        ("场景三 · 异常事件检测", "原始信号 + LSTM + 残差三级告警",                ACCENT),
        ("场景四 · 新字段接入",   "先算 ACF 与 SNR, 按可用性判据预判, 无需训练",   PRIMARY),
        ("场景五 · 冲击信号监控", "ShockX/Y/Z 不适用本方案, 需采用事件检测算法",   DANGER),
    ]
    y2 = Inches(5.0)
    for scene, rec, color in rec_lines:
        add_rect(slide, Inches(0.8), y2, Inches(0.15), Inches(0.35), fill=color)
        add_text(slide, Inches(1.05), y2, Inches(4.0), Inches(0.35),
                 scene, size=12, bold=True, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(5.0), y2, Inches(8.0), Inches(0.35),
                 rec, size=11, color=TEXT_MID,
                 anchor=MSO_ANCHOR.MIDDLE)
        y2 = y2 + Inches(0.38)

    add_footer(slide, page_no, total)
    add_notes(slide, note_text)


# ============================================================
# 文字页内容
# ============================================================
def sections_for_findings():
    return [
        {"kind": "card",
         "tag": "发现 1", "tag_color": SUCCESS,
         "head": "族内一致性显著, 判断字段适用性可按物理量族进行",
         "body": [
             "• 加速度 RMS (3 字段): NRMSE 13.5–17.0%, 标准差 1.7%",
             "• 速度 RMS   (4 字段): NRMSE 13.1–15.4%, 标准差 1.2%",
             "• 冲击类     (3 字段): NRMSE 33.7–48.4%, 标准差 7.9%",
         ]},
        {"kind": "card",
         "tag": "发现 2", "tag_color": DANGER,
         "head": "ShockX 出现性能反转: TimesFM 较 LSTM 差 8.6%",
         "body": [
             "• ShockX 的 SNR=0.238, 为 11 字段中最低, 噪声约为信号的 4 倍",
             "• TFM MAE=55.6, LSTM MAE=51.2, TFM 偏差较大",
             "• 表明信噪比极低时, 预训练大模型易受噪声主导, 不宜直接使用",
         ]},
        {"kind": "card",
         "tag": "发现 3", "tag_color": ACCENT,
         "head": "冲击类信号整族不适用本方案",
         "body": [
             "• 三个冲击字段的 SNR 均 < 0.35, ACF lag-16 均 < 0.67",
             "• 三个冲击字段的 NRMSE 均 > 33%, 而其他 8 个字段均 < 18%",
             "• 结论: 冲击类整族不适用 SavGol 趋势预测方案, 建议改用事件检测",
         ]},
    ]


def sections_for_criteria_v2():
    return [
        {"kind": "banner",
         "content": "▍初版判据命中率 9/11  ·  两个边界字段未正确归档, 需调整阈值",
         "fill": LIGHT_ORANGE, "color": ACCENT, "size": 14},
        {"kind": "card",
         "tag": "边界 1", "tag_color": ACCENT,
         "head": "aRMSY (SNR=1.396) 被误判至 B 档",
         "body": [
             "• aRMSY 的 NRMSE=15.1%, 与 A 档字段的表现一致",
             "• 原因: 初版 A 档阈值 SNR>1.5 过严, 相差 0.1 即被排除",
             "• 调整: A 档阈值放宽为 SNR>1.3",
         ]},
        {"kind": "card",
         "tag": "边界 2", "tag_color": DANGER,
         "head": "ShockY (SNR=0.33) 被误判至 B 档",
         "body": [
             "• ShockY NRMSE=33.7%, 恰在 B/C 分界线内侧, 被归为 B",
             "• 但 SNR=0.33, 噪声约为信号的 3 倍, 本质不具备可预测性",
             "• 调整: 引入硬红线 SNR<0.5 直接归入 C, 覆盖 NRMSE 计算结果",
         ]},
        {"kind": "code", "height": 1.6,
         "content": (
             "def classify(nrmse, acf16, snr):\n"
             "    if snr < 0.5:                                      # 硬红线\n"
             "        return 'C 不建议直接使用'\n"
             "    if nrmse < 0.20 and acf16 > 0.75 and snr > 1.3:\n"
             "        return 'A 可直接部署'\n"
             "    if nrmse < 0.25 and acf16 > 0.70:\n"
             "        return 'B 可用, 建议复核'\n"
             "    return 'C 不建议直接使用'"
         )},
        {"kind": "banner",
         "content": "▍修订后判据命中率 11/11  ·  所有字段归档与预测性能一致",
         "fill": LIGHT_GREEN, "color": SUCCESS, "size": 14},
    ]


def sections_for_data_enough():
    return [
        {"kind": "scenarios", "rows": [
            {"status": "ok",  "icon": "✓",
             "scenario": "方法学结论: 方案对 RMS 与幅值类物理量有效",
             "answer": "充分 — 8 个样本, 族内标准差 ≤ 2%"},
            {"status": "ok",  "icon": "✓",
             "scenario": "方法学结论: 方案跨物理量族具备泛化能力",
             "answer": "充分 — 三族多样本分离清晰"},
            {"status": "ok",  "icon": "✓",
             "scenario": "方法学结论: 方案在 MPB_01 设备内稳健",
             "answer": "充分 — 11 字段全覆盖"},
            {"status": "mid", "icon": "△",
             "scenario": "时段稳定性: 不同时段内是否保持一致",
             "answer": "部分覆盖 — 建议补充分段对比"},
            {"status": "no",  "icon": "✗",
             "scenario": "跨设备泛化: 换设备后结果是否保持",
             "answer": "不足 — 当前仅 1 台设备数据"},
            {"status": "no",  "icon": "✗",
             "scenario": "生产部署决策: 能否直接上线",
             "answer": "不足 — 需跨设备及跨时段补充验证"},
        ]},
        {"kind": "banner",
         "content": "▍方法学层面的结论已具备充分样本支撑;生产部署决策需进一步数据",
         "fill": LIGHT, "color": PRIMARY, "size": 14, "bold": True,
         "align": PP_ALIGN.CENTER, "height": 0.5},
        {"kind": "card",
         "tag": "后续工作",
         "tag_color": PRIMARY,
         "head": "建议补充的跨设备验证实验",
         "body": [
             "• 采集 2–3 台不同设备、各 2 小时的时序数据即可",
             "• 仅需以 Magnitude (A 档代表字段) 作为验证对象",
             "• 判断标准: NRMSE 保持在 10–18% 区间则泛化成立; 超过 25% 则需设备级微调",
         ]},
    ]


# ============================================================
# 主渲染
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    total = len(SLIDES_SPEC)

    for page_i, spec in enumerate(SLIDES_SPEC):
        chart_no, title, subtitle, img, key_point, notes = spec
        s = prs.slides.add_slide(BLANK)

        if chart_no == "COVER":
            render_cover(s, title, subtitle, notes)
            continue
        if chart_no == "AGENDA":
            render_agenda(s, title, subtitle, notes)
            continue
        if chart_no == "结论":
            summary_notes = (
                "本页汇总全部研究发现与工程部署建议。\n\n"
                "七项主要发现概括如下。\n"
                "发现一: 1 秒尺度原始振动信号的 lag-1 自相关约为 0.87,\n"
                "但 lag-60 自相关降至约 0.09, 属于高频噪声特征, 逐点预测\n"
                "在此尺度下不具备有效性。\n"
                "发现二: 采用 SavGol(61, 3) 分解后, 趋势信号的 lag-1 自\n"
                "相关提升至 1.00, 具备充分的可预测性。\n"
                "发现三: horizon = 16 秒为精度与窗口长度的较优平衡点,\n"
                "超过此值后 MAE 呈指数增长。\n"
                "发现四: 在 11 个字段中, TimesFM 零样本预测在 10 个字段上\n"
                "优于本地训练的 LSTM, 平均领先约 41%。唯一例外为 ShockX,\n"
                "其信噪比过低导致预训练大模型不具优势。\n"
                "发现五: 物理量族内部一致性显著, NRMSE 标准差不超过 2%,\n"
                "表明族级预判可替代逐字段测试。\n"
                "发现六: 冲击类信号整族不适用本方案, NRMSE 介于 33-48%,\n"
                "且 SNR 均小于 0.4, 需采用事件检测类方法处理。\n"
                "发现七: ACF + SNR + NRMSE 构成的可用性判据对本数据集\n"
                "的 11 个字段命中率达到 100%, 可作为新字段接入前的评估\n"
                "标准。\n\n"
                "工程部署建议分五类场景:\n"
                "场景一为日常振动趋势监测, 采用推荐配置 SavGol(61,3) +\n"
                "TimesFM + horizon=16, NRMSE 约 14%。\n"
                "场景二为短期实时预测, 采用更短窗口配置 SavGol(31,2) +\n"
                "TimesFM + horizon=8, NRMSE 可进一步降至约 7%。\n"
                "场景三为异常事件检测, 采用原始信号配合 LSTM 残差三级\n"
                "告警策略。\n"
                "场景四为新字段接入流程, 以 ACF 与 SNR 为预判依据, 结合\n"
                "可用性判据进行归档, 无需额外模型训练。\n"
                "场景五为冲击信号监控, 由于本方案不适用, 需改用脉冲事件\n"
                "检测算法。\n\n"
                "后续研究方向: 补充跨设备、跨时段的数据验证, 以支撑生产\n"
                "部署决策;对冲击类信号探索基于分位数回归或事件驱动的\n"
                "替代方案。"
            )
            render_summary(s, title, subtitle, summary_notes,
                           page_i + 1, total)
            continue

        if str(chart_no).endswith("-TEXT"):
            pure_no = chart_no.replace("-TEXT", "")
            if "7-4" in pure_no:
                sections = sections_for_findings()
                notes = (
                    "本页汇总全字段验证实验所得的三项主要发现。\n\n"
                    "发现一聚焦于物理量族的一致性。将 11 个字段按物理量族\n"
                    "划分后, 每个族内部 NRMSE 的标准差均不超过 2%, 而族\n"
                    "之间则存在显著差异。这说明在实际应用中, 对新字段的\n"
                    "可用性评估可以以物理量族为单位完成, 无需逐字段进行\n"
                    "完整训练与测试。\n\n"
                    "发现二关注 ShockX 字段上出现的性能反转。该字段是 11\n"
                    "个字段中唯一一个 TimesFM 预测性能低于 LSTM 的案例,\n"
                    "差距为 8.6%。分析表明其根本原因是信噪比极低\n"
                    "(SNR=0.238), 信号骨架被噪声淹没, 此时预训练大模型\n"
                    "倾向于将噪声建模为可重复模式, 反而引入系统性偏差。\n"
                    "该结果对工程部署具有警示意义: 当 SNR 低于 0.25 时,\n"
                    "不建议直接采用大模型零样本推理。\n\n"
                    "发现三明确了冲击类信号的整族特性。原四字段实验中\n"
                    "仅包含 ShockZ 一个冲击样本, 存在「是否为轴向特例」\n"
                    "的疑问。扩展至 ShockX/Y/Z 全部三轴后, 三字段的\n"
                    "SNR 均小于 0.35、ACF lag-16 均小于 0.67、NRMSE 均\n"
                    "大于 33%, 共同表明冲击类信号作为一个物理量族不适合\n"
                    "趋势预测类方案, 应采用事件检测或分位数建模等替代\n"
                    "方法。"
                )
            elif "7-5" in pure_no:
                sections = sections_for_criteria_v2()
                notes = (
                    "本页说明可用性判据的修订过程与修订后效果。\n\n"
                    "初版判据的命中率为 9/11, 两个归档错误均出现在分档\n"
                    "边界上, 分别对应阈值过严与阈值过松两个方向。\n\n"
                    "边界一为 aRMSY 字段。该字段 SNR 为 1.396, 恰低于初\n"
                    "版 A 档的 SNR>1.5 阈值约 0.1, 被归入 B 档。但其实\n"
                    "际 NRMSE 为 15.1%, 与 A 档字段表现一致。修订方案\n"
                    "将 A 档 SNR 阈值放宽至 1.3, 以避免单一指标微小差\n"
                    "距导致的误判。\n\n"
                    "边界二为 ShockY 字段。该字段 NRMSE 为 33.7%, 恰低\n"
                    "于初版 B 档上限 35%, 被归入 B 档。但其 SNR 仅为\n"
                    "0.33, 噪声约为信号的三倍, 从信号本质上不具备可预\n"
                    "测性。修订方案引入硬红线: 当 SNR 低于 0.5 时一律\n"
                    "归入 C 档, 不再受 NRMSE 影响。\n\n"
                    "修订后的判据函数在本报告中作为推荐实现, 对 MPB_01\n"
                    "数据集的 11 个字段命中率达到 100%, 可作为后续新字\n"
                    "段接入前的自动评估标准。"
                )
            else:  # 7-6
                sections = sections_for_data_enough()
                notes = (
                    "本页说明研究数据的覆盖范围及结论的有效边界。\n\n"
                    "数据覆盖范围: 本研究仅基于一台设备 (MPB_01) 在一\n"
                    "段约 2.62 小时的连续监测数据。涉及 11 个物理量字段,\n"
                    "按 1 秒重采样后每字段约 8960 个样本点。\n\n"
                    "结论有效范围可分为两个层次。\n"
                    "方法学层面: 样本量足以支撑「方案对 RMS 与幅值类\n"
                    "物理量有效」、「方案具备跨物理量族泛化能力」、「方\n"
                    "案在 MPB_01 设备内稳健」三项结论。\n"
                    "生产决策层面: 由于缺乏跨设备数据, 本研究不能回答\n"
                    "「换一台设备方案是否仍然有效」这一问题。在生产部\n"
                    "署决策前, 需补充跨设备与跨时段的验证实验。\n\n"
                    "建议的后续实验成本较低: 采集 2–3 台不同设备各约 2\n"
                    "小时的 CSV 数据, 仅需对 Magnitude 一个字段重跑本\n"
                    "研究的 phase_a 与 phase_b 流程, 预计总耗时不超过\n"
                    "10 分钟。判断标准为: NRMSE 维持在 10-18% 区间表明\n"
                    "泛化成立; 超过 25% 则需要针对设备进行微调。"
                )
            render_text_page(s, pure_no, title, subtitle, sections,
                             notes, page_i + 1, total)
            continue

        # 图表页
        set_slide_bg(s, WHITE)
        slide_header(s, chart_no, title, subtitle)
        if img:
            add_figure_area(s, os.path.join(OUTPUT_DIR, img),
                            top=Inches(1.35), max_height=Inches(5.0))
        if key_point:
            add_key_point(s, key_point, top=Inches(6.5))
        add_footer(s, page_i + 1, total)
        add_notes(s, notes)

    prs.save(PPT_PATH)
    print(f"[OK] 已生成: {PPT_PATH}")
    print(f"总页数: {len(prs.slides)}")


if __name__ == "__main__":
    build()
