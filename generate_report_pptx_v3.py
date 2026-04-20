"""
生成项目报告 PPT (v3)
基于 v2,在"章节 6 多属性泛化"之后追加"章节 7 全字段验证":
  - 新的 NRMSE 指标说明 (修正旧的 MAE/|均值|)
  - 11 字段 NRMSE 对比图
  - 判据象限图 (ACF × SNR)
  - 3 大发现文字页
  - 判据 v2 (命中率 11/11) 文字页
  - 数据够不够? 分场景文字页
  - 更新总结页
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT_DIR = r"E:\timesfm_project\output"
PPT_PATH = os.path.join(OUTPUT_DIR, "振动预测项目报告_v3.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# 明亮色调(沿用 v2)
PRIMARY = RGBColor(0x4A, 0x90, 0xE2)    # 天蓝
ACCENT = RGBColor(0xF5, 0xA6, 0x23)     # 暖橙
SUCCESS = RGBColor(0x7E, 0xD3, 0x21)    # 草绿
DANGER = RGBColor(0xE7, 0x4C, 0x3C)     # 红
BG = RGBColor(0xFC, 0xFD, 0xFE)
LIGHT = RGBColor(0xE8, 0xF0, 0xFA)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_GREEN = RGBColor(0xEB, 0xF8, 0xD6)
LIGHT_RED = RGBColor(0xFD, 0xEC, 0xEA)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_MID = RGBColor(0x66, 0x6E, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_font(run, name="微软雅黑", size=14, bold=False, color=TEXT_DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for e in rPr.findall(qn("a:ea")):
        rPr.remove(e)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)


def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.text = ""
        run = p.add_run()
        run.text = line
        set_font(run, name=font_name, size=size, bold=bold, color=color)
    return tb


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, note_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = ""
    lines = note_text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line


def slide_header(slide, chart_no, chart_name, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(0.3), Inches(1.8), Inches(0.45),
             fill=ACCENT)
    add_text(slide, Inches(0.4), Inches(0.32), Inches(1.8), Inches(0.4),
             chart_no, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(2.4), Inches(0.3), Inches(10.5), Inches(0.5),
             chart_name, size=20, bold=True, color=PRIMARY)
    if subtitle:
        add_text(slide, Inches(2.4), Inches(0.78), Inches(10.5), Inches(0.3),
                 subtitle, size=11, color=TEXT_MID)


def add_figure_area(slide, img_path, top=Inches(1.35), max_height=Inches(5.2)):
    if not os.path.exists(img_path):
        return
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    max_w = SLIDE_W - Inches(0.8)
    ratio = iw / ih
    if max_height * ratio > max_w:
        w = max_w
        h = w / ratio
    else:
        h = max_height
        w = h * ratio
    left = (SLIDE_W - w) // 2
    slide.shapes.add_picture(img_path, left, top, width=w, height=h)
    return top + h


def add_key_point(slide, text, top=None):
    if top is None:
        top = Inches(6.7)
    add_rect(slide, Inches(0.4), top, SLIDE_W - Inches(0.8), Inches(0.55),
             fill=LIGHT)
    add_text(slide, Inches(0.6), top + Inches(0.05),
             SLIDE_W - Inches(1.2), Inches(0.45),
             "▍核心观点  " + text, size=13, bold=True, color=PRIMARY,
             anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_no, total):
    add_text(slide, Inches(11.8), SLIDE_H - Inches(0.3),
             Inches(1.5), Inches(0.3),
             f"{page_no} / {total}",
             size=10, color=TEXT_MID, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.3),
             Inches(11), Inches(0.3),
             "MPB_01 振动预测与预防性维护研究",
             size=10, color=TEXT_MID)


# ============================================================
# 新增: 纯文字页面渲染 (用于发现/判据/够不够)
# ============================================================
def render_text_page(slide, chart_no, title, subtitle, sections, note_text,
                     page_no, total):
    """
    sections: list of dict:
      { "kind": "card",  "tag": "发现 1", "tag_color": PRIMARY,
        "head": "族内一致性非常稳",
        "body": ["• ...", "• ..."]
      }
      { "kind": "code", "content": "def ..." }
      { "kind": "table", "rows": [[...], [...]], "headers": [...] }
    """
    set_slide_bg(slide, WHITE)
    slide_header(slide, chart_no, title, subtitle)

    y = Inches(1.35)
    for sec in sections:
        if sec["kind"] == "card":
            # 左侧标签
            tag_color = sec.get("tag_color", PRIMARY)
            add_rect(slide, Inches(0.5), y, Inches(1.3), Inches(1.0),
                     fill=tag_color)
            add_text(slide, Inches(0.5), y, Inches(1.3), Inches(1.0),
                     sec["tag"], size=15, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 右侧主体
            body_h = Inches(1.0)
            add_rect(slide, Inches(1.8), y, SLIDE_W - Inches(2.3), body_h,
                     fill=LIGHT)
            add_text(slide, Inches(2.0), y + Inches(0.06),
                     SLIDE_W - Inches(2.7), Inches(0.4),
                     sec["head"], size=14, bold=True, color=PRIMARY)
            body_text = "\n".join(sec["body"])
            add_text(slide, Inches(2.0), y + Inches(0.42),
                     SLIDE_W - Inches(2.7), Inches(0.6),
                     body_text, size=11, color=TEXT_DARK)
            y = y + body_h + Inches(0.12)
        elif sec["kind"] == "code":
            h = Inches(sec.get("height", 1.6))
            add_rect(slide, Inches(0.5), y, SLIDE_W - Inches(1.0), h,
                     fill=RGBColor(0x2C, 0x3E, 0x50))
            add_text(slide, Inches(0.7), y + Inches(0.1),
                     SLIDE_W - Inches(1.4), h - Inches(0.2),
                     sec["content"], size=12, color=WHITE,
                     font_name="Consolas")
            y = y + h + Inches(0.1)
        elif sec["kind"] == "scenarios":
            # 分场景表, 左场景 + 右结论(颜色 icon)
            for s in sec["rows"]:
                row_h = Inches(0.45)
                color = {"ok": SUCCESS, "no": DANGER, "mid": ACCENT}[s["status"]]
                bg_col = {"ok": LIGHT_GREEN, "no": LIGHT_RED, "mid": LIGHT_ORANGE}[s["status"]]
                add_rect(slide, Inches(0.5), y, Inches(0.6), row_h, fill=color)
                add_text(slide, Inches(0.5), y, Inches(0.6), row_h,
                         s["icon"], size=14, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                add_rect(slide, Inches(1.1), y, SLIDE_W - Inches(1.6), row_h,
                         fill=bg_col)
                add_text(slide, Inches(1.3), y, Inches(8.0), row_h,
                         s["scenario"], size=12, bold=True, color=TEXT_DARK,
                         anchor=MSO_ANCHOR.MIDDLE)
                add_text(slide, Inches(9.5), y, SLIDE_W - Inches(10.2), row_h,
                         s["answer"], size=11, color=TEXT_MID,
                         anchor=MSO_ANCHOR.MIDDLE)
                y = y + row_h + Inches(0.06)
        elif sec["kind"] == "banner":
            # 大横幅文字
            add_rect(slide, Inches(0.5), y, SLIDE_W - Inches(1.0),
                     Inches(sec.get("height", 0.6)),
                     fill=sec.get("fill", LIGHT))
            add_text(slide, Inches(0.7), y, SLIDE_W - Inches(1.4),
                     Inches(sec.get("height", 0.6)),
                     sec["content"], size=sec.get("size", 14),
                     bold=sec.get("bold", True),
                     color=sec.get("color", PRIMARY),
                     align=sec.get("align", PP_ALIGN.LEFT),
                     anchor=MSO_ANCHOR.MIDDLE)
            y = y + Inches(sec.get("height", 0.6)) + Inches(0.1)

    add_footer(slide, page_no, total)
    add_notes(slide, note_text)


# ============================================================
# Spec (v3: v2 原内容 + 新章节 7)
# ============================================================
SLIDES_SPEC = [
    ("COVER", "MPB_01 设备振动预测与预防性维护研究",
     "LSTM vs TimesFM · 信号分解 · 多参数对比 · 11 字段泛化性验证",
     None, None,
     "本次汇报覆盖 5 大主题(v3 版):\n"
     "  1. 11 属性原始数据的可视化与统计分析\n"
     "  2. LSTM 与 TimesFM 两种算法在不同参数下的对比\n"
     "  3. 为什么预测结果会退化为直线——诊断分析\n"
     "  4. 信号分解+预测窗口优化后的有效方案\n"
     "  5. 全字段泛化验证与可用性判据 (v3 新增)\n"
     "数据来源: MPB01_6.11 08_18.csv (InfluxDB 导出)\n"
     "数据规模: 30 万行, 11 属性, ~2.6 小时连续监测\n"
     "设备: 1 台 MPB_01 (注意: 未验证跨设备泛化)"),

    ("AGENDA", "报告目录",
     "本次汇报分 6 个章节 (v3 版)",
     None, None,
     "章节 1: 数据概览与可视化\n"
     "章节 2: 两种算法多参数对比\n"
     "章节 3: 预测直线问题诊断\n"
     "章节 4: 信号分解与窗口优化\n"
     "章节 5: 多属性泛化测试 (4 字段)\n"
     "章节 6: 全字段验证与可用性判据 (v3 新增, 11 字段)\n"
     "章节 7: 项目结论与工程建议"),

    ("图 1-1", "11 属性分量级时序对照",
     "按数值量级分 3 组, 橙色带标注 18 次疑似停机",
     "属性总览_1_分量级时序.png",
     "Magnitude 存在 18 次骤降事件, 对应设备停机或传感器离线",
     "图表内容:\n"
     "  - 【组A】Magnitude (合成幅值, 0-1 量级)\n"
     "  - 【组B】aRMS/vRMS 系列 (均方根, 0-5 量级)\n"
     "  - 【组C】ShockX/Y/Z (三轴冲击, ±1500 量级)\n"
     "  - 橙色半透明带 = 检测到的 18 次疑似异常事件\n\n"
     "参数/术语:\n"
     "  - RMS (Root Mean Square) 均方根: 衡量信号能量\n"
     "  - Magnitude: 三轴合成后的总幅值\n"
     "  - aRMS (acceleration RMS) 加速度均方根, 单位 g\n"
     "  - vRMS (velocity RMS) 速度均方根, 单位 mm/s\n"
     "  - Shock: 冲击峰值, 反映瞬时撞击强度\n"
     "  - 重采样频率 = 1 秒, 每个子图有 8960 个点"),

    ("图 1-2", "11 属性独立时序图",
     "每个属性一个面板, 便于观察各自量级和波动",
     "属性总览_2_独立子图.png",
     "Shock 类是尖峰脉冲信号, RMS 类是连续起伏, 物理本质不同",
     "图表内容:\n"
     "  - 11 个属性各自独立面板, 以便突出各自特点\n"
     "  - 每个子图标题含数据范围和标准差\n"
     "  - 红色虚线为均值线, 便于观察偏离\n\n"
     "参数/术语:\n"
     "  - 均值 (mean): 所有值的算术平均\n"
     "  - STD (Standard Deviation) 标准差: 衡量偏离均值程度\n"
     "  - 范围 [min, max]: 数据上下限\n\n"
     "观察:\n"
     "  - ShockX/Y/Z: 密集尖峰(脉冲型)\n"
     "  - aRMS/vRMS: 随时间起伏(连续型)\n"
     "  - Magnitude: 平稳段 + 多次骤降(停机)"),

    ("图 1-3", "多属性统计概览与相关性",
     "log 轴量级对比 + 变异系数 + 11×11 相关矩阵",
     "属性总览_3_统计与相关性.png",
     "Magnitude/aRMS/vRMS 本质同一物理量, Shock 独立",
     "图表内容:\n"
     "  左上: 各属性绝对值均值 log 柱状图 (揭示跨量级差异)\n"
     "  右上: 变异系数 CV, ShockY=1.66 最难预测\n"
     "  下方: 11×11 皮尔逊相关系数热力图\n\n"
     "参数/术语:\n"
     "  - CV (变异系数) = STD / |均值|\n"
     "    CV > 1   = 高难度 (红色线)\n"
     "    CV > 0.5 = 中难度 (橙色线)\n"
     "  - 皮尔逊相关系数 r ∈ [-1, 1]\n"
     "    |r| > 0.7 强相关; |r| < 0.3 几乎无关\n\n"
     "关键发现:\n"
     "  - Magnitude 与 aRMS/vRMS 相关 0.86-0.98 (本质同类)\n"
     "  - Shock 与其他独立 (相关 ~0.05)"),

    ("图 1-4", "异常事件放大 — Magnitude 骤降",
     "18 次事件的前 4 个详细波形",
     "属性总览_4_异常事件放大.png",
     "典型事件: 振动从 ~0.86 骤降到 ~0, 持续 3-14 秒后恢复",
     "图表内容:\n"
     "  展示 18 次疑似异常事件中的前 4 个\n"
     "  每个子图放大显示骤降前后 30 秒\n"
     "  橙色阴影 = 异常持续区间\n"
     "  灰色虚线 = 正常中位数\n\n"
     "参数/术语:\n"
     "  - z_threshold = 3.0 (异常判定阈值, 越低越敏感)\n"
     "  - min_duration = 3 秒 (最短持续, 避免瞬时误判)\n"
     "  - 判定规则: 值 < 中位数 - 3.0 * STD * 0.3\n\n"
     "业务含义:\n"
     "  - 骤降到 0 很可能是设备停机或传感器离线\n"
     "  - 对预测评估的启示: 测试集若含异常段会干扰评估"),

    ("图 2-1", "LSTM vs TimesFM 分维度参数对比",
     "在 Magnitude 上测试 3 个维度 × 7 组实验的 MAE",
     "对比_1_分维度柱状图.png",
     "TimesFM 在绝大多数配置下都优于 LSTM",
     "图表内容:\n"
     "  横向 3 个面板, 分别展示 3 个参数维度:\n"
     "  - 维度1 重采样频率: 1s / 5s / 30s\n"
     "  - 维度2 预测长度 horizon: 8 / 32 / 64\n"
     "  - 维度3 上下文长度 context: 128 / 256 / 512\n"
     "  红柱 = LSTM, 绿柱 = TimesFM, 越低越好\n\n"
     "参数/术语:\n"
     "  - MAE (Mean Absolute Error) 平均绝对误差\n"
     "  - horizon: 一次预测未来几步\n"
     "  - context: 模型看多长的历史数据\n\n"
     "结果:\n"
     "  - 5s 重采样比 1s 稍好 (降噪)\n"
     "  - horizon 越小 MAE 越小\n"
     "  - TimesFM 对 context 不敏感, LSTM 对 context 敏感"),

    ("图 2-2", "6 组参数下的预测曲线样例",
     "发现问题: 预测普遍退化为近似直线",
     "对比_2_预测曲线样例.png",
     "两种模型的预测曲线都偏平, 触发后续诊断",
     "图表内容:\n"
     "  6 个子图对应 6 组不同参数实验\n"
     "  深蓝 = 真实值, 红虚线 = LSTM, 绿点划线 = TimesFM\n"
     "  可见: 即便是最好的 horizon=8, 模型曲线仍显著平于\n"
     "  真实波动\n\n"
     "参数/术语:\n"
     "  - 真实值: 测试集的实际振动值\n"
     "  - LSTM: Encoder-Decoder + 直接多步输出, 本地训练\n"
     "  - TimesFM: Google 预训练基础模型 2.5-200M\n\n"
     "启示:\n"
     "  - 表象 MAE 不差, 但视觉跟随效果不佳\n"
     "  - 需要后续诊断找出根本原因"),

    ("图 2-3", "算法胜率与 MAE 相对差距汇总",
     "TimesFM 在 5/6 组中胜出 (83%), 平均领先 41%",
     "对比_3_胜率汇总.png",
     "TimesFM 零样本即用, LSTM 需训练但效果不稳定",
     "图表内容:\n"
     "  左图: 胜率饼图, TimesFM 83% vs LSTM 17%\n"
     "  右图: 每组实验的 MAE 相对差距横向条形\n\n"
     "参数/术语:\n"
     "  - MAE 差距% = (差的/好的 - 1) × 100%\n"
     "  - 例如 LSTM=0.112, TFM=0.079, 差距 = 42.9%\n\n"
     "时间成本:\n"
     "  - TimesFM 总耗时 71 秒 (仅推理)\n"
     "  - LSTM 总耗时 137 秒 (训练+推理)"),

    ("图 3-1", "诊断: 为什么预测结果是直线?",
     "自相关 + 朴素基线 + 功率谱三维度分析",
     "诊断_1_可预测性分析.png",
     "1 秒尺度接近高频噪声, 最优预测本就是均值线",
     "图表内容:\n"
     "  左上: 不同重采样尺度下的自相关函数 (ACF)\n"
     "  右上: 1s 信号功率谱 PSD\n"
     "  中间: 朴素基线 vs LSTM/TimesFM MAE 对比\n"
     "  左下: 各尺度的相对误差\n"
     "  右下: 100 秒真实振动切片 (高频抖动)\n\n"
     "参数/术语:\n"
     "  - ACF (AutoCorrelation Function) 自相关函数\n"
     "    r(lag=k) = x[t] 与 x[t+k] 的相关性\n"
     "    r > 0.3 才有预测价值 (红色阈值线)\n"
     "  - 朴素基线: 不用模型的简单方法\n"
     "    例 朴素-最后值: 未来每步都等于当前最后值\n"
     "  - PSD (Power Spectral Density) 功率谱\n\n"
     "关键发现:\n"
     "  - 1s lag-1 自相关 0.87 (相邻秒高度相关)\n"
     "  - 但 lag-60 只有 0.09 (远期几乎独立)\n"
     "  - 朴素-最后值 MAE=0.069, 比 LSTM (0.112) 还好"),

    ("图 4-1", "方案 B - 信号分解对比",
     "把 Magnitude 拆成「趋势 + 噪声」, 测试 3 种方法",
     "分解_1_信号分解对比.png",
     "SavGol 保留趋势形状最好, Butter 最平滑但损失细节",
     "图表内容:\n"
     "  第一行: 原始信号(灰) + 3 种趋势(彩色)叠加\n"
     "  后 3 行: 每种分解的「趋势(主)+残差(副轴)」\n\n"
     "参数/术语:\n"
     "  - 移动均值 (Moving Average): 窗口内取平均\n"
     "    window = 30 秒\n"
     "  - Savitzky-Golay 滤波: 局部多项式拟合\n"
     "    window_length = 61, polyorder = 3\n"
     "  - Butterworth 低通滤波: 频域滤波\n"
     "    cutoff_hz = 0.05, order = 4\n"
     "  - 趋势占比 = Var(趋势) / Var(原信号)\n\n"
     "结果:\n"
     "  - 移动均值 70.5% / SavGol 71.9% / Butter 84.1%"),

    ("图 4-2", "方案 B - 分解后可预测性跃升",
     "ACF 从 0.87 提升到 1.00, MAE 下降 47%",
     "分解_2_可预测性对比.png",
     "分解后的趋势信号真正可预测, 模型不再只能预测均值",
     "图表内容:\n"
     "  左图: ACF 对比 (lag-1 蓝, lag-30 紫)\n"
     "  右图: MAE 对比 (朴素灰, LSTM 红, TimesFM 绿)\n\n"
     "数据对比:\n"
     "  原始信号:  lag-1=0.87, LSTM=0.111, TFM=0.084\n"
     "  移动均值:  lag-1=1.00, LSTM=0.067, TFM=0.038\n"
     "  SavGol ★: lag-1=1.00, LSTM=0.055, TFM=0.044\n"
     "  Butter:   lag-1=1.00, LSTM=0.107, TFM=0.060"),

    ("图 4-3", "方案 B - 分解后预测曲线样例",
     "SavGol 趋势行: 预测曲线终于能跟上波动",
     "分解_3_预测曲线对比.png",
     "信号分解后预测视觉效果明显改善",
     "图表内容:\n"
     "  4 行 × 3 列, 每行一种信号, 每列一个样例\n"
     "  行1 原始 / 行2 移动均值 / 行3 SavGol / 行4 Butter\n\n"
     "观察:\n"
     "  - 原始信号行: 真实值高频抖动, 预测平均线\n"
     "  - SavGol 行: 预测曲线弯曲跟随 ✓\n"
     "  - horizon 仍是 32, 所以仍有改进空间"),

    ("图 5-1", "缩短预测窗口 horizon 的视觉效果",
     "horizon = 4 / 8 / 16 / 32 在 SavGol 趋势上",
     "horizon_1_缩短窗口对比.png",
     "horizon=16 是最佳平衡: 预测曲线明显跟随, MAE 低",
     "图表内容:\n"
     "  4 行 × 4 列, 每行一个 horizon, 每列一个样例\n"
     "  可见:\n"
     "    h=4  太短看不出\n"
     "    h=8  TimesFM 开始跟随\n"
     "    h=16 ★ 两个模型都完美跟随真实弧度\n"
     "    h=32 又退化为直线\n\n"
     "参数/术语:\n"
     "  - horizon: 预测未来多少秒\n"
     "  - context: 固定 256 秒历史\n"
     "  - 信号: SavGol(61,3) 分解的趋势部分"),

    ("图 5-2", "预测窗口与精度的定量关系",
     "MAE 随 horizon 指数增长, TimesFM 始终领先",
     "horizon_2_MAE曲线.png",
     "horizon = 16 相对误差仅 2.1%, 推荐工程使用",
     "图表内容:\n"
     "  MAE vs horizon 折线图\n"
     "  红色 = LSTM, 绿色 = TimesFM\n\n"
     "精度对照:\n"
     "  h=4  TFM MAE=0.005 (相对误差 0.5%)\n"
     "  h=8  TFM MAE=0.009 (相对误差 1.0%)\n"
     "  h=16 TFM MAE=0.018 (相对误差 2.1%) ← 推荐\n"
     "  h=32 TFM MAE=0.051 (相对误差 5.9%)\n\n"
     "注意: 这里的相对误差是 MAE/|均值|,\n"
     "后续在章节 6 我们会修正为更科学的 NRMSE"),

    ("图 6-1", "多属性泛化性 — 预测曲线",
     "最佳配置 (SavGol + h=16) 在 4 个代表属性上",
     "多属性_1_预测曲线.png",
     "连续类 3 属性效果佳, 脉冲类 ShockZ 不适用",
     "图表内容:\n"
     "  4 行 × 3 列, 每行一个属性, 每列一个样例\n"
     "  属性顺序: Magnitude / aRMSX / vRMSM / ShockZ\n\n"
     "观察:\n"
     "  - Magnitude/aRMSX/vRMSM: 预测明显跟随真实值 ✓\n"
     "  - ShockZ: 预测偏平 (脉冲型不适合此方案)"),

    ("图 6-2", "多属性泛化性 — 指标汇总 (旧指标)",
     "TimesFM 在连续类属性上相对误差稳定 ~2.1% (存在陷阱)",
     "多属性_2_指标对比.png",
     "旧指标: ShockZ 看上去也是 2.4%, 其实是陷阱 → 见章节 6",
     "图表内容:\n"
     "  左: MAE 对比 (log 轴)\n"
     "  中: 相对误差% 对比\n"
     "  右: 趋势信号 ACF 对比\n\n"
     "【注意: 此页使用的是「相对误差 = MAE / |均值|」】\n"
     "ShockZ 均值 |−226| 非常大, 分母把 MAE 稀释成 2.4%\n"
     "看上去和其他属性一样好, 但实际上预测是失败的\n"
     "我们在章节 6 会用 NRMSE 替代此指标, 暴露真实水平"),

    # ========== 章节 6 新增 ==========
    ("图 7-1", "指标修正 — 从「MAE/|均值|」到 NRMSE",
     "旧指标会被大均值稀释, NRMSE 更能反映真实预测难度",
     "多属性_3_NRMSE对比.png",
     "换指标后 ShockZ 的真实 NRMSE 是 36%, 不是 2.4%",
     "【为什么要换指标】\n"
     "  旧指标: 相对误差% = MAE / |均值| × 100%\n"
     "  问题: 当均值很大(如 ShockZ 均值 -226),\n"
     "        这个分母会稀释掉真实误差, 造成「看着很好」的假象\n\n"
     "【新指标 NRMSE】\n"
     "  NRMSE = MAE / 趋势std × 100%\n"
     "  含义: 预测误差相对于「信号本身的波动幅度」的比例\n"
     "        NRMSE=10% 表示预测误差只有信号波动的 1/10\n"
     "        NRMSE=100% 表示完全无效(跟预测均值一样)\n\n"
     "【对比效果 (4 字段)】\n"
     "  旧指标 Magnitude/aRMSX/vRMSM/ShockZ:\n"
     "    2.1% / 2.2% / 2.1% / 2.4%  → 看起来都一样好\n"
     "  新指标 NRMSE:\n"
     "    13.8% / 13.5% / 13.1% / 36.1%  → ShockZ 真实很差\n\n"
     "【结论】\n"
     "  所有后续报告应统一使用 NRMSE\n"
     "  绿色 20% 线为「可用」标准, 红色 35% 为「不建议」边界"),

    ("图 7-2", "全字段验证 — 11 字段 NRMSE 对比",
     "把 4 字段扩展到 11 字段, 按物理量族分色",
     "全字段_1_NRMSE对比.png",
     "8 个 RMS/幅值字段 NRMSE 13-17%, 3 个冲击字段 33-48%",
     "【这张图干什么】\n"
     "  把原 4 字段 (Magnitude/aRMSX/vRMSM/ShockZ) 扩展到\n"
     "  所有 11 个字段, 验证「A/B/C 可用性判据」是否成立\n\n"
     "【如何读图】\n"
     "  - 斜线填充 = LSTM, 实心 = TimesFM\n"
     "  - 颜色编码物理量族:\n"
     "    蓝 = 合成幅值, 青绿 = 加速度RMS, 紫 = 速度RMS, 红 = 冲击类\n"
     "  - 绿色虚线 20% = A 档可直接上线\n"
     "  - 红色虚线 35% = C 档不建议\n\n"
     "【族内一致性 (最重要的观察)】\n"
     "  加速度RMS (3 字段): NRMSE 13.5% - 17.0%, std=1.7%\n"
     "  速度RMS (4 字段):   NRMSE 13.1% - 15.4%, std=1.2%\n"
     "  冲击类 (3 字段):    NRMSE 33.7% - 48.4%, std=7.9%\n\n"
     "  解读: 同一物理量族内方差极小\n"
     "        → 判断新字段能不能用, 看族就够了, 不需逐字段测试\n\n"
     "【TimesFM 全胜 LSTM (除 ShockX)】\n"
     "  10/11 字段 TFM 优于 LSTM, 平均领先 41%\n"
     "  唯一例外: ShockX (TFM 比 LSTM 还差 8.6%)\n"
     "  说明 SNR<0.25 时预训练大模型反而有害"),

    ("图 7-3", "全字段验证 — 可用性判据象限图",
     "ACF × SNR 两个预算定标准判断信号能不能用",
     "全字段_2_判据象限.png",
     "冲击 3 点全部聚集在左下(难预测区), RMS 8 点全部在右上",
     "【这张图的核心价值】\n"
     "  一眼看懂「什么样的信号能预测、什么样的不能」\n"
     "  两个坐标轴 + 气泡大小 = 三维信息, 完全不需文字\n\n"
     "【坐标轴含义】\n"
     "  X 轴: ACF lag-16 (0~1)\n"
     "    意思: 16 秒前的值能否影响当前的值\n"
     "    >0.75 → 信号有「记忆」, 可预测\n"
     "    <0.60 → 信号基本「忘了」过去, 难预测\n"
     "  Y 轴: SNR (对数) = 趋势std / 噪声std\n"
     "    意思: 信号的「骨架」比「毛刺」大多少倍\n"
     "    >1.5 → 规律远大于噪声\n"
     "    <0.5 → 噪声淹没规律\n\n"
     "【气泡含义】\n"
     "  颜色 = 物理量族\n"
     "  大小 = NRMSE 大小(大=预测误差大)\n\n"
     "【物理分离】\n"
     "  右上角(ACF>0.8, SNR>1.4): 8 个 RMS/幅值字段\n"
     "  左下角(ACF<0.7, SNR<0.4): 3 个冲击字段\n"
     "  两个区域中间完全没有过渡 → 判据非常干净\n\n"
     "【业务用法】\n"
     "  新接入一个字段时, 只需算 ACF 和 SNR 两个数就能\n"
     "  预判效果, 不用真的训练模型"),

    ("图 7-4-TEXT", "全字段验证 — 3 大发现",
     "族内稳定 + ShockX 灾难 + 冲击类整族失败",
     None, None, None),  # 用 TEXT 类型渲染

    ("图 7-5-TEXT", "可用性判据 v2 — 命中率从 9/11 → 11/11",
     "两处微调: SNR 阈值放宽到 1.3, 加 SNR<0.5 硬红线",
     None, None, None),

    ("图 7-6-TEXT", "数据够不够? — 分场景回答",
     "方法学结论已足够, 跨设备泛化仍需更多数据",
     None, None, None),

    ("总结", "项目结论与工程建议 (v3 更新)",
     "7 项关键发现 + 5 种业务场景的推荐配置",
     None, None, None),
]


def render_cover(slide, title, subtitle, note_text):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.25), fill=PRIMARY)
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), fill=ACCENT)

    add_text(slide, Inches(0.5), Inches(2.3), SLIDE_W - Inches(1), Inches(1.4),
             title, size=40, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5), Inches(3.9), Inches(3.3), Inches(0.04),
             fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(4.1), SLIDE_W - Inches(1), Inches(0.5),
             subtitle, size=18, color=TEXT_MID, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(5.3), SLIDE_W - Inches(1), Inches(0.4),
             "数据: MPB01_6.11 08_18.csv  ·  设备: MPB_01  ·  11 属性时序",
             size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(5.7), SLIDE_W - Inches(1), Inches(0.4),
             "v3 · 扩展至 11 字段全集, 引入 NRMSE 指标与 A/B/C 可用性判据",
             size=12, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    add_notes(slide, note_text)


def render_agenda(slide, title, subtitle, note_text):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.35), SLIDE_W - Inches(1), Inches(0.7),
             title, size=32, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(1.1), SLIDE_W - Inches(1), Inches(0.4),
             subtitle, size=14, color=TEXT_MID, align=PP_ALIGN.CENTER)

    chapters = [
        ("01", "数据概览与可视化", "11 属性统计与时序分析", PRIMARY),
        ("02", "两种算法多参数对比", "LSTM vs TimesFM 网格实验", ACCENT),
        ("03", "预测直线问题诊断", "ACF 与朴素基线分析", SUCCESS),
        ("04", "信号分解与窗口优化", "SavGol 滤波 + 短 horizon 方案", PRIMARY),
        ("05", "多属性泛化(4 字段)", "初步泛化观察", ACCENT),
        ("06", "全字段验证(11 字段) ★ v3 新增", "NRMSE 指标 + A/B/C 判据", DANGER),
        ("07", "结论与工程建议", "判据 v2 + 5 场景推荐", SUCCESS),
    ]
    y = Inches(1.7)
    for num, ctitle, csub, color in chapters:
        add_rect(slide, Inches(1.5), y, Inches(1.0), Inches(0.68),
                 fill=color)
        add_text(slide, Inches(1.5), y, Inches(1.0), Inches(0.68),
                 num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(2.5), y, Inches(9.5), Inches(0.68),
                 fill=LIGHT)
        add_text(slide, Inches(2.7), y + Inches(0.05),
                 Inches(9.3), Inches(0.35),
                 ctitle, size=14, bold=True, color=PRIMARY)
        add_text(slide, Inches(2.7), y + Inches(0.38),
                 Inches(9.3), Inches(0.3),
                 csub, size=10, color=TEXT_MID)
        y = y + Inches(0.78)

    add_notes(slide, note_text)


def render_summary(slide, title, subtitle, note_text, page_no, total):
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.25), SLIDE_W - Inches(1), Inches(0.5),
             title, size=26, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.75), SLIDE_W - Inches(1), Inches(0.3),
             subtitle, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

    findings = [
        "1s 尺度原始振动接近高频噪声, 逐点预测必然退化",
        "SavGol 信号分解能有效分离可预测的「趋势」部分",
        "horizon = 16 秒是预测精度与业务价值的最佳平衡",
        "TimesFM 零样本胜出 LSTM: 10/11 字段领先, 平均提升 41%",
        "物理量族内一致性极强(std ≤ 2%), 可按族判断而非逐字段",
        "冲击类信号整族不适合趋势预测 (NRMSE 33-48%, SNR<0.4)",
        "可用性判据(ACF + SNR + NRMSE)命中率 11/11, 可作为上线前的 checklist",
    ]

    y = Inches(1.2)
    for i, f in enumerate(findings):
        add_rect(slide, Inches(0.5), y, Inches(1.2), Inches(0.38),
                 fill=PRIMARY if i % 2 == 0 else ACCENT)
        add_text(slide, Inches(0.5), y, Inches(1.2), Inches(0.38),
                 f"发现{i+1}", size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(1.7), y, Inches(11.0), Inches(0.38),
                 fill=LIGHT)
        add_text(slide, Inches(1.9), y,
                 Inches(10.8), Inches(0.38),
                 f, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.44)

    add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
             fill=LIGHT)
    add_text(slide, Inches(0.8), Inches(4.55),
             Inches(12), Inches(0.4),
             "▍工程部署推荐配置",
             size=15, bold=True, color=PRIMARY)

    rec_lines = [
        ("场景 1 · 振动趋势监测",  "SavGol(61,3) + TimesFM + h=16, NRMSE ~14%",   SUCCESS),
        ("场景 2 · 短期实时预测",  "SavGol(31,2) + TimesFM + h=8,  NRMSE ~7%",    SUCCESS),
        ("场景 3 · 异常事件检测",  "原始信号 + LSTM + 残差 3 级告警",              ACCENT),
        ("场景 4 · 新字段接入",    "先算 ACF + SNR, 用 A/B/C 判据预判, 不用训练",  PRIMARY),
        ("场景 5 · 冲击信号监控",  "ShockX/Y/Z 不用此方案, 改用脉冲事件检测",      DANGER),
    ]
    y2 = Inches(5.0)
    for scene, rec, color in rec_lines:
        add_rect(slide, Inches(0.8), y2, Inches(0.15), Inches(0.35), fill=color)
        add_text(slide, Inches(1.05), y2, Inches(4.0), Inches(0.35),
                 scene, size=12, bold=True, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(5.0), y2, Inches(8.0), Inches(0.35),
                 rec, size=11, color=TEXT_MID,
                 anchor=MSO_ANCHOR.MIDDLE)
        y2 = y2 + Inches(0.38)

    add_footer(slide, page_no, total)
    add_notes(slide, note_text)


# ============================================================
# 纯文字页的具体内容
# ============================================================
def sections_for_findings():
    return [
        {"kind": "card",
         "tag": "发现 1", "tag_color": SUCCESS,
         "head": "族内一致性非常稳 — 按族判断即可, 不用逐字段测",
         "body": [
             "• 加速度RMS 族 (3 字段): NRMSE 13.5–17.0%, 标准差 1.7%",
             "• 速度RMS   族 (4 字段): NRMSE 13.1–15.4%, 标准差 1.2%",
             "• 冲击类   族 (3 字段): NRMSE 33.7–48.4%, 标准差 7.9%",
         ]},
        {"kind": "card",
         "tag": "发现 2", "tag_color": DANGER,
         "head": "ShockX 是灾难: TFM 输给 LSTM 8.6% (唯一反例)",
         "body": [
             "• ShockX: SNR=0.238 (噪声≈信号 4 倍), 11 字段里最差",
             "• TFM MAE=55.6, LSTM MAE=51.2  →  大模型把噪声学进去了",
             "• 结论: SNR<0.25 时预训练大模型不仅没用, 反而有害",
         ]},
        {"kind": "card",
         "tag": "发现 3", "tag_color": ACCENT,
         "head": "冲击类不是 Z 轴特例, 是整族失败",
         "body": [
             "• 3 个冲击字段 SNR 全 < 0.35, ACF_lag16 全 < 0.67",
             "• 3 个冲击字段 NRMSE 全 > 33% (其他 8 字段都 < 18%)",
             "• 结论: 冲击族整体不适合 SavGol 趋势预测 → 需换方案",
         ]},
    ]


def sections_for_criteria_v2():
    return [
        {"kind": "banner",
         "content": "▍判据 v1 (原版) 命中率: 9/11 — 两个边界案例归错档",
         "fill": LIGHT_ORANGE, "color": ACCENT, "size": 14},
        {"kind": "card",
         "tag": "问题 1", "tag_color": ACCENT,
         "head": "A 档 SNR>1.5 太严格, aRMSY (SNR=1.396) 被踢到 B",
         "body": [
             "• aRMSY 的 NRMSE=15.1%, 其实跟 A 档均值完全一样好",
             "• 原因: SNR 差 0.1 就被规则一刀切",
             "• 修复: SNR>1.5 放宽到 SNR>1.3",
         ]},
        {"kind": "card",
         "tag": "问题 2", "tag_color": DANGER,
         "head": "B 档过宽, ShockY (SNR=0.33) 误归 B, 实际不能用",
         "body": [
             "• ShockY NRMSE=33.7%, 踩在 B/C 线上被归 B",
             "• 但 SNR=0.33 (噪声 3 倍于信号), 跟 ShockZ 本质一样",
             "• 修复: 加硬红线 SNR<0.5 一律 C (无论 NRMSE 多少)",
         ]},
        {"kind": "code", "height": 1.6,
         "content": (
             "def classify_v2(nrmse, acf16, snr):\n"
             "    if snr < 0.5:                          # 硬红线\n"
             "        return 'C 不建议直接用'\n"
             "    if nrmse < 0.20 and acf16 > 0.75 and snr > 1.3:\n"
             "        return 'A 可直接上线'\n"
             "    if nrmse < 0.25 and acf16 > 0.70:\n"
             "        return 'B 可用/建议复核'\n"
             "    return 'C 不建议直接用'"
         )},
        {"kind": "banner",
         "content": "▍判据 v2 命中率: 11/11 — 所有字段归档正确",
         "fill": LIGHT_GREEN, "color": SUCCESS, "size": 14},
    ]


def sections_for_data_enough():
    return [
        {"kind": "scenarios", "rows": [
            {"status": "ok",  "icon": "✓",
             "scenario": "方法学结论: 方案对 RMS/幅值类物理量有效",
             "answer": "够 — 8 个样本,族内 std ≤ 2%"},
            {"status": "ok",  "icon": "✓",
             "scenario": "方法学结论: 方案跨物理量族有泛化能力",
             "answer": "够 — 3 族 × 多样本,分离清晰"},
            {"status": "ok",  "icon": "✓",
             "scenario": "方法学结论: 方案在 MPB_01 设备上 robust",
             "answer": "够 — 11 字段全覆盖"},
            {"status": "mid", "icon": "△",
             "scenario": "设备稳定性: 时段内是否稳定",
             "answer": "部分够 — 可切 3 段再测 (建议做)"},
            {"status": "no",  "icon": "✗",
             "scenario": "跨设备泛化: 换一台机器还能用吗?",
             "answer": "不够 — 需要 2–3 台机器 CSV"},
            {"status": "no",  "icon": "✗",
             "scenario": "生产上线决策",
             "answer": "不够 — 需跨设备 + 跨时段验证"},
        ]},
        {"kind": "banner",
         "content": "▍结论  方法学论文 / 技术选型报告: 够了 ·  生产部署决策: 还需 2 台机器数据",
         "fill": LIGHT, "color": PRIMARY, "size": 14, "bold": True,
         "align": PP_ALIGN.CENTER, "height": 0.5},
        {"kind": "card",
         "tag": "下一步",
         "tag_color": PRIMARY,
         "head": "如果要做生产决策, 请尽快获取跨设备数据",
         "body": [
             "• 需要 2-3 台不同机器, 各 2 小时 csv 即可",
             "• 只挑 Magnitude 一个字段跑 (代表 A 档, 速度快)",
             "• 预期结果: NRMSE 仍在 10-18% → 方案过关; 超过 25% → 需设备级 fine-tune",
         ]},
    ]


# ============================================================
# 主渲染
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    # 需要在生成前就知道总页数
    total = len(SLIDES_SPEC)

    for page_i, spec in enumerate(SLIDES_SPEC):
        chart_no, title, subtitle, img, key_point, notes = spec
        s = prs.slides.add_slide(BLANK)

        if chart_no == "COVER":
            render_cover(s, title, subtitle, notes)
            continue
        if chart_no == "AGENDA":
            render_agenda(s, title, subtitle, notes)
            continue
        if chart_no == "总结":
            # 新的总结用自定义 note
            summary_notes = (
                "【v3 总结页 — 请逐条讲】\n"
                "发现 1: 原始振动 1s 级相邻相关 0.87 但 60s 相关仅 0.09,\n"
                "        属于高频噪声, 逐点预测必然退化\n"
                "发现 2: SavGol 61,3 分解后 lag-1 ACF 提升到 1.00,\n"
                "        趋势信号真正可预测\n"
                "发现 3: horizon=16 是 MAE 和业务价值的拐点, 超过就爆炸\n"
                "发现 4: 11 字段里 10 个 TFM 胜 LSTM, 平均提升 41%\n"
                "        唯一例外是 ShockX (SNR 太低, 大模型失效)\n"
                "发现 5: 物理量族是最强的泛化单位, 族内 std ≤ 2%\n"
                "        说明: 同族内任何一个字段测过一次, 其他默认同档\n"
                "发现 6: 冲击类整族 NRMSE 33-48%, 不适合趋势预测,\n"
                "        需改用事件驱动方法\n"
                "发现 7: ACF + SNR + NRMSE 三指标判据 v2 命中率 11/11,\n"
                "        可直接作为新字段上线前的 checklist\n\n"
                "【5 个场景】\n"
                "场景 1: 日常趋势监测, 最常用的生产配置\n"
                "场景 2: 短期实时预警, 滚动窗口可达 1% 级误差\n"
                "场景 3: 异常检测用残差+告警, 不用 TFM\n"
                "场景 4: 接入新字段时, 不训练模型, 只算 ACF/SNR 判档\n"
                "场景 5: 冲击类别走此方案, 必须换事件检测算法\n\n"
                "【下一步】\n"
                "最重要的一件事: 拿 2-3 台不同机器的 CSV 重跑验证\n"
                "这是目前结论最薄弱的一环, 不补上不能做生产决策"
            )
            render_summary(s, title, subtitle, summary_notes,
                           page_i + 1, total)
            continue

        # 新类型: TEXT 页 (纯文字页面)
        if str(chart_no).endswith("-TEXT"):
            pure_no = chart_no.replace("-TEXT", "")
            if "7-4" in pure_no:
                sections = sections_for_findings()
                notes = (
                    "【3 大发现解读】\n\n"
                    "发现 1 · 族内一致性 (最有商业价值的结论)\n"
                    "  意思: 把 11 个字段按物理量分成 4 族后,\n"
                    "        同族内部的 NRMSE 标准差只有 1-2%\n"
                    "  为什么重要: 未来新设备接入时,\n"
                    "        不用逐字段跑测试, 看族就能预判效果\n"
                    "  数字: 加速度族 std=1.7%, 速度族 std=1.2%\n\n"
                    "发现 2 · ShockX 灾难 (最应该拿出来讲的警告)\n"
                    "  意思: ShockX 是 11 字段里唯一一个 TFM 输给 LSTM 的\n"
                    "  为什么重要: 证明「预训练大模型不是万能的」\n"
                    "        信号质量差到一定程度时, 大模型反而有害\n"
                    "  数字: SNR=0.238, TFM MAE 55.6 > LSTM MAE 51.2 (-8.6%)\n\n"
                    "发现 3 · 冲击类整族失败 (回答质疑的关键证据)\n"
                    "  意思: 原来只测 ShockZ 时别人可能问「是不是 Z 轴特殊」\n"
                    "        现在 X/Y/Z 全测, 证明是整族问题\n"
                    "  数字: 3 个字段 SNR 全 < 0.35, NRMSE 全 > 33%\n"
                    "  行动建议: 冲击类不用预测模型, 用事件检测"
                )
            elif "7-5" in pure_no:
                sections = sections_for_criteria_v2()
                notes = (
                    "【判据 v2 解读】\n\n"
                    "【原判据(v1)为什么不够好】\n"
                    "  v1 的 A 档要求 SNR>1.5 太严格\n"
                    "  aRMSY 的 SNR=1.396 差一点点就被踢去 B\n"
                    "  但它的 NRMSE=15.1% 实际跟 A 档字段完全一样好\n"
                    "  → 这不是信号不行, 是判据太机械\n\n"
                    "  v1 的 B 档范围太宽\n"
                    "  ShockY NRMSE=33.7%, 勉强没超过 35% 被归了 B\n"
                    "  但它的 SNR=0.33, 噪声比信号大 3 倍\n"
                    "  → 这种信号不管 NRMSE 算出来多少都不能用\n\n"
                    "【修复策略 (2 条改动)】\n"
                    "  修复 1: SNR 阈值从 1.5 降到 1.3 (放宽 A 档)\n"
                    "  修复 2: 加一条 SNR<0.5 的硬红线 (收紧 B 档)\n\n"
                    "【修复后的效果】\n"
                    "  aRMSY: SNR=1.396 > 1.3  ✓  进 A\n"
                    "  ShockY: SNR=0.33 < 0.5  ✓  进 C (硬红线)\n"
                    "  命中率 9/11 → 11/11\n\n"
                    "【这段代码的用法】\n"
                    "  以后任何新字段接入, 拿 3 个数字(NRMSE/ACF/SNR)\n"
                    "  跑一下这个函数, 返回档位就直接决策\n"
                    "  不用训练模型, 秒级决策"
                )
            else:  # 7-6
                sections = sections_for_data_enough()
                notes = (
                    "【数据够不够的分场景解读】\n\n"
                    "【必须诚实: 数据只来自 1 台机器 × 1 段时间】\n"
                    "  原始 CSV: MPB01_6.11 08_18.csv\n"
                    "  时长 2.62 小时 (不是之前以为的 10 小时)\n"
                    "  11 个物理量字段, 单机器 MPB_01\n\n"
                    "【分场景回答「够不够」】\n"
                    "  方法学论文: 够 — 11 字段 × 2 算法已经很硬\n"
                    "  技术选型报告: 够 — 判据命中 11/11\n"
                    "  设备稳定性: 部分够 — 可切 3 段再测 (成本低)\n"
                    "  跨设备泛化: 不够 — 根本没有第二台机器数据\n"
                    "  生产部署决策: 不够 — 必须补跨设备验证\n\n"
                    "【讲这页时最重要的话】\n"
                    "  「方法学上我们已经把能做的都做了,\n"
                    "   但要推到生产, 必须拿到 2-3 台机器的数据再跑一遍」\n\n"
                    "【下一步成本】\n"
                    "  2-3 台机器, 每台 2 小时 CSV\n"
                    "  只跑 Magnitude 一个字段 (A 档代表)\n"
                    "  总耗时 < 10 分钟\n"
                    "  判决标准: NRMSE 10-18% 过关, >25% 需 fine-tune"
                )
            render_text_page(s, pure_no, title, subtitle, sections,
                             notes, page_i + 1, total)
            continue

        # 普通图表页
        set_slide_bg(s, WHITE)
        slide_header(s, chart_no, title, subtitle)
        if img:
            add_figure_area(s, os.path.join(OUTPUT_DIR, img),
                            top=Inches(1.35), max_height=Inches(5.0))
        if key_point:
            add_key_point(s, key_point, top=Inches(6.5))
        add_footer(s, page_i + 1, total)
        add_notes(s, notes)

    prs.save(PPT_PATH)
    print(f"[OK] 已生成: {PPT_PATH}")
    print(f"总页数: {len(prs.slides)}")


if __name__ == "__main__":
    build()
